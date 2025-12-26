#!/usr/bin/env python3
"""
Bolagsverket MCP Server v5.1.0
==============================
Hämtar och analyserar företagsdata från Bolagsverkets API "Värdefulla datamängder".

NYA FUNKTIONER I v5.1.0:
- K3-koncernredovisning: Stöd för koncernnyckeltal och K3K-taxonomin
- Taxonomiversions-detektion: Automatisk identifiering och varning för arkiverade versioner
- BAS-kontoreferenser: Mappning mellan taxonomi-begrepp och BAS-kontoplanen
- Revisionsberättelse-parsing: Strukturerad extraktion av revisionsberättelse
- Decimals-precision: Behåller full precision från iXBRL
- Fastställelseintyg: Extraktion av beslut och signaturer
- Caching med TTL: SQLite-baserad cache för bättre prestanda
- Utökad information: Stöd för extension taxonomy (odefinierade begrepp, rubrikändringar)

ÄNDRINGAR I v5.0.1:
- Fixat: sign-attribut hanteras nu korrekt (taxonomier.se Exempel 1d)
- Fixat: format-attribut hanteras nu korrekt (taxonomier.se Exempel 1e)

NYA FUNKTIONER I v5.0:
- Pydantic field_validator för robust organisationsnummer-validering
- Utökade Resources med fler URI-endpoints
- Nya Prompts för vanliga arbetsflöden
- Förbättrad felhantering med strukturerade felmeddelanden

ALLA VERKTYG:
1. bolagsverket_analyze_full - Komplett årsredovisningsanalys
2. bolagsverket_koncern - Koncernanalys med K3K-stöd (NYTT i v5.1)
3. bolagsverket_risk_check - Röda flaggor och varningar  
4. bolagsverket_export_pro - Export till Word/PowerPoint
5. bolagsverket_trend - Historisk tidsserie & prognoser
6. bolagsverket_search - Sök företag utan orgnummer
7. bolagsverket_network - Styrelsenätverk & personkopplingar
8. bolagsverket_taxonomy_info - Taxonomiversion och varningar (NYTT i v5.1)
9. bolagsverket_bas_mapping - BAS-kontomappning (NYTT i v5.1)
10. bolagsverket_cache_stats - Cache-statistik (NYTT i v5.1)

MCP Best Practices Implementation:
- Tools: Granulära verktyg för specifika operationer
- Resources: Passiv data via URI-schema (utökat i v5.1)
- Prompts: Fördefinierade arbetsflöden (utökat i v5.0)
- Pydantic: Field validators för input-validering (nytt i v5.0)
- iXBRL: Korrekt sign/format-attribut-hantering (fixat i v5.0.1)
- Caching: SQLite med TTL (nytt i v5.1)
- Logging: stderr (aldrig stdout för STDIO-transport)
- Strukturerad felhantering
"""

import json
import uuid
import httpx
import zipfile
import sys
import os
import re
import csv
import logging
import sqlite3
import hashlib
from datetime import datetime, timedelta
from typing import Optional, Dict, Any, List, Tuple, Set
from enum import Enum
from io import BytesIO, StringIO
from dataclasses import dataclass, asdict, field
from pathlib import Path
from collections import defaultdict
import statistics

from bs4 import BeautifulSoup
try:
    from fastmcp import FastMCP
except ImportError:
    from mcp.server.fastmcp import FastMCP
from pydantic import BaseModel, Field, ConfigDict, field_validator

# =============================================================================
# FÖRBÄTTRING #1: Logging till stderr (ALDRIG stdout för STDIO-transport)
# =============================================================================
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    stream=sys.stderr  # KRITISKT: Måste vara stderr för MCP STDIO
)
logger = logging.getLogger("bolagsverket_mcp")

# Valfria imports för export
try:
    import openpyxl
    from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    logger.warning("openpyxl ej installerat - Excel-export inaktiverad")

try:
    from weasyprint import HTML, CSS
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logger.warning("weasyprint ej installerat - PDF-export inaktiverad")


# =============================================================================
# Konfiguration - Läser från miljövariabler med fallback
# =============================================================================

CLIENT_ID = os.environ.get("BOLAGSVERKET_CLIENT_ID", "UIiATHgXGSP6HIyOlqWZkX51dnka")
CLIENT_SECRET = os.environ.get("BOLAGSVERKET_CLIENT_SECRET", "H10hBNr_KeYqA9h5AEe7J32HkFsa")
TOKEN_URL = os.environ.get("BOLAGSVERKET_TOKEN_URL", "https://portal.api.bolagsverket.se/oauth2/token")
BASE_URL = os.environ.get("BOLAGSVERKET_BASE_URL", "https://gw.api.bolagsverket.se/vardefulla-datamangder/v1")
SCOPE = "vardefulla-datamangder:read vardefulla-datamangder:ping"

# Output-katalog: använd /tmp på server, Downloads lokalt
OUTPUT_DIR = Path(os.environ.get("OUTPUT_DIR", "/tmp/bolagsverket" if os.environ.get("RENDER") else str(Path.home() / "Downloads" / "bolagsverket")))


# =============================================================================
# FÖRBÄTTRING #10: Kodlistor enligt API-dokumentationen
# =============================================================================

# Avregistreringsorsak (The reason an organisation was removed from the register)
AVREGISTRERINGSORSAK = {
    "AKEJH": "Aktiekapitalet inte höjts",
    "ARSEED": "Årsredovisning saknas",
    "AVREG": "Avregistrerad",
    "BABAKEJH": "Ombildat till bankaktiebolag eller aktiekapitalet inte höjts",
    "DELAV": "Delning",
    "DOM": "Beslut av domstol",
    "FUAV": "Fusion",
    "GROMAV": "Gränsöverskridande ombildning",
    "KKAV": "Konkurs",
    "LIAV": "Likvidation",
    "NYINN": "Ny innehavare",
    "OMAV": "Ombildning",
    "OMBAB": "Ombildat till bankaktiebolag",
    "OVERK": "Overksamhet",
    "UTLKKLI": "Det utländska företagets likvidation eller konkurs",
    "VERKUPP": "Verksamheten har upphört",
    "VDSAK": "Verkställande direktör saknas",
}

# Organisationsform (Form of organisation)
ORGANISATIONSFORM = {
    "AB": "Aktiebolag",
    "BAB": "Bankaktiebolag",
    "BF": "Bostadsförening",
    "BFL": "Utländsk banks filial",
    "BRF": "Bostadsrättsförening",
    "E": "Enskild näringsverksamhet",
    "EB": "Enkla bolag",
    "EEIG": "Europeisk ekonomisk intressegruppering",
    "EGTS": "Europeiska grupperingar för territoriellt samarbete",
    "EK": "Ekonomisk förening",
    "FAB": "Försäkringsaktiebolag",
    "FF": "Försäkringsförmedlare",
    "FL": "Filial",
    "FOF": "Försäkringsförening",
    "HB": "Handelsbolag",
    "I": "Ideell förening som bedriver näringsverksamhet",
    "KB": "Kommanditbolag",
    "KHF": "Kooperativ hyresrättsförening",
    "MB": "Medlemsbank",
    "OFB": "Ömsesidigt försäkringsbolag",
    "OTPB": "Ömsesidigt tjänstepensionsbolag",
    "S": "Stiftelse som bedriver näringsverksamhet",
    "SB": "Sparbank",
    "SCE": "Europakooperativ",
    "SE": "Europabolag",
    "SF": "Sambruksförening",
    "TPAB": "Tjänstepensionsaktiebolag",
    "TPF": "Tjänstepensionsförening",
    "TSF": "Trossamfund som bedriver näringsverksamhet",
}

# Pågående avvecklings- eller omstruktureringsförfarande
PAGAENDE_FORFARANDE = {
    "AC": "Ackordsförhandling",
    "DEOL": "Överlåtande vid delning",
    "DEOT": "Övertagande vid delning",
    "FR": "Företagsrekonstruktion",
    "FUOL": "Överlåtande i fusion",
    "FUOT": "Övertagande i fusion",
    "GROM": "Gränsöverskridande ombildning",
    "KK": "Konkurs",
    "LI": "Likvidation",
    "OM": "Ombildning",
    "RES": "Resolution",
}

# Identitetsbeteckningstyp för organisation
IDENTITETSBETECKNINGSTYP_ORG = {
    "DODSBO": "Dödsbonummer",
    "GDNUMMER": "Identitetsbeteckning person (GD-nummer)",
    "ORGANISATIONSNUMMER": "Organisationsnummer",
    "PERSONNUMMER": "Identitetsbeteckning person (personnummer)",
    "SAMORDNINGSNUMMER": "Identitetsbeteckning person (Samordningsnummer)",
    "UTLANDSK_JURIDISK_IDENTITETSBETECKNING": "Utländsk identitetsbeteckning",
}

# Organisationsnamntyp
ORGANISATIONSNAMNTYP = {
    "FORETAGSNAMN": "Företagsnamn",
    "FORNAMN_FRSPRAK": "Företagsnamn på främmande språk",
    "NAMN": "Namn",
    "SARSKILT_FORETAGSNAMN": "Särskilt företagsnamn",
}


def get_avregistreringsorsak_text(kod: str) -> str:
    """Hämta klartext för avregistreringsorsak med fallback."""
    return AVREGISTRERINGSORSAK.get(kod, kod)


def get_organisationsform_text(kod: str) -> str:
    """Hämta klartext för organisationsform med fallback."""
    return ORGANISATIONSFORM.get(kod, kod)


def get_pagaende_forfarande_text(kod: str) -> str:
    """Hämta klartext för pågående förfarande med fallback."""
    return PAGAENDE_FORFARANDE.get(kod, kod)


# =============================================================================
# FÖRBÄTTRING #10: Strukturerade felkoder
# =============================================================================

class ErrorCode(str, Enum):
    COMPANY_NOT_FOUND = "COMPANY_NOT_FOUND"
    ANNUAL_REPORT_NOT_FOUND = "ANNUAL_REPORT_NOT_FOUND"
    API_ERROR = "API_ERROR"
    AUTH_ERROR = "AUTH_ERROR"
    PARSE_ERROR = "PARSE_ERROR"
    INVALID_INPUT = "INVALID_INPUT"
    EXPORT_ERROR = "EXPORT_ERROR"


@dataclass
class MCPError:
    """Strukturerat fel enligt MCP best practices."""
    code: ErrorCode
    message: str
    details: Dict[str, Any] = field(default_factory=dict)
    
    def to_response(self) -> str:
        """Formatera som MCP-kompatibelt felmeddelande."""
        return json.dumps({
            "isError": True,
            "errorCode": self.code.value,
            "message": self.message,
            "details": self.details
        }, ensure_ascii=False, indent=2)


def handle_error(code: ErrorCode, message: str, **details) -> str:
    """Skapa strukturerat felmeddelande."""
    error = MCPError(code=code, message=message, details=details)
    logger.error(f"{code.value}: {message} - {details}")
    return error.to_response()


# =============================================================================
# Dataklasser
# =============================================================================

@dataclass
class Person:
    fornamn: str
    efternamn: str
    roll: str
    datum: Optional[str] = None
    
    @property
    def fullnamn(self) -> str:
        return f"{self.fornamn} {self.efternamn}".strip()


@dataclass
class Nyckeltal:
    nettoomsattning: Optional[int] = None
    resultat_efter_finansiella: Optional[int] = None
    arets_resultat: Optional[int] = None
    eget_kapital: Optional[int] = None
    balansomslutning: Optional[int] = None
    soliditet: Optional[float] = None
    antal_anstallda: Optional[int] = None
    vinstmarginal: Optional[float] = None
    roe: Optional[float] = None
    
    def berakna_nyckeltal(self):
        if self.nettoomsattning and self.arets_resultat:
            self.vinstmarginal = round((self.arets_resultat / self.nettoomsattning) * 100, 2)
        if self.eget_kapital and self.arets_resultat and self.eget_kapital > 0:
            self.roe = round((self.arets_resultat / self.eget_kapital) * 100, 2)


@dataclass
class Arsredovisning:
    org_nummer: str
    foretag_namn: str
    rakenskapsar_start: str
    rakenskapsar_slut: str
    nyckeltal: Nyckeltal
    personer: List[Person]
    balansrakning: Dict[str, Any]
    resultatrakning: Dict[str, Any]
    noter: Dict[str, str]
    metadata: Dict[str, str]


@dataclass 
class CompanyInfo:
    """Grundläggande företagsinformation."""
    org_nummer: str
    namn: str
    organisationsform: str
    juridisk_form: Optional[str]
    registreringsdatum: str
    status: str
    avregistreringsdatum: Optional[str]
    avregistreringsorsak: Optional[str]
    adress: Dict[str, str]
    verksamhet: Optional[str]
    sni_koder: List[Dict[str, str]]
    sate: Optional[str]
    # Nya fält för riskanalys
    pagaende_konkurs: Optional[Dict[str, str]] = None  # {datum, typ}
    pagaende_likvidation: Optional[Dict[str, str]] = None  # {datum, typ}
    reklamsparr: Optional[bool] = None
    verksam_organisation: Optional[bool] = None
    # Nya fält från API-specifikationen (förbättring #7)
    registreringsland: Optional[Dict[str, str]] = None  # {kod, klartext}
    namnskyddslopnummer: Optional[int] = None  # Viktigt för enskilda firmor
    infört_hos_scb: Optional[str] = None  # Datum när org infördes hos SCB
    organisationsform_kod: Optional[str] = None  # AB, HB, KB, E etc.
    alla_namn: Optional[List[Dict[str, Any]]] = None  # Alla registrerade namn
    datakalla_fel: Optional[List[Dict[str, str]]] = None  # Fel från datakällor


@dataclass
class RodFlagga:
    """Varning/röd flagga för ett företag."""
    typ: str
    allvarlighet: str  # "kritisk", "varning", "info"
    beskrivning: str
    varde: Optional[Any] = None
    rekommendation: Optional[str] = None


@dataclass
class FullArsredovisning:
    """Komplett årsredovisning med ALL information."""
    org_nummer: str
    foretag_namn: str
    rakenskapsar_start: str
    rakenskapsar_slut: str
    nyckeltal: Nyckeltal
    styrelse: List[Person]
    revisorer: List[Person]
    vd: Optional[Person]
    balansrakning: Dict[str, Any]
    resultatrakning: Dict[str, Any]
    forvaltningsberattelse: Dict[str, str]
    noter: Dict[str, str]
    flerarsdata: List[Dict[str, Any]]
    roda_flaggor: List[RodFlagga]
    metadata: Dict[str, str]


@dataclass
class Trendanalys:
    """Trendanalys för ett företag över tid."""
    org_nummer: str
    foretag_namn: str
    perioder: List[str]
    nyckeltal_serie: Dict[str, List[Optional[int]]]
    tillvaxt: Dict[str, Optional[float]]
    prognos: Dict[str, Optional[int]]


@dataclass
class PersonNetwork:
    """Person med alla deras bolagsengagemang."""
    namn: str
    bolag: List[Dict[str, Any]]


# =============================================================================
# NYA DATAKLASSER FÖR v5.1.0
# =============================================================================

@dataclass
class TaxonomiInfo:
    """Information om använd taxonomi (FÖRBÄTTRING 2)."""
    version: str
    typ: str  # K2, K3, K3K, REVISION, FASTSTALLELSE
    entry_point: str
    ar_arkiverad: bool = False
    varning: Optional[str] = None

@dataclass
class KoncernNyckeltal:
    """Nyckeltal för koncernredovisning (FÖRBÄTTRING 1)."""
    koncern_nettoomsattning: Optional[int] = None
    koncern_rorelseresultat: Optional[int] = None
    koncern_resultat_efter_finansiella: Optional[int] = None
    koncern_arets_resultat: Optional[int] = None
    koncern_eget_kapital: Optional[int] = None
    koncern_balansomslutning: Optional[int] = None
    minoritetsandel: Optional[int] = None
    goodwill: Optional[int] = None
    koncern_soliditet: Optional[float] = None
    
    def berakna_koncern_nyckeltal(self):
        """Beräkna härledda koncernnyckeltal."""
        if self.koncern_eget_kapital and self.koncern_balansomslutning:
            self.koncern_soliditet = round(
                (self.koncern_eget_kapital / self.koncern_balansomslutning) * 100, 1
            )

@dataclass
class Revisionsberattelse:
    """Strukturerad revisionsberättelse (FÖRBÄTTRING 4)."""
    revisor_namn: Optional[str] = None
    revisor_titel: Optional[str] = None
    revisionsbolag: Optional[str] = None
    uttalande_arsredovisning: Optional[str] = None
    uttalande_koncernredovisning: Optional[str] = None
    uttalande_forvaltning: Optional[str] = None
    grund_for_uttalande: Optional[str] = None
    anmarkningar: List[str] = field(default_factory=list)
    ovrigt: Optional[str] = None
    datum: Optional[str] = None
    ort: Optional[str] = None
    ar_ren: bool = True  # Ren revisionsberättelse utan anmärkningar
    typ: str = "standard"  # standard, koncern, forenklad

@dataclass
class PrecisionValue:
    """Värde med precision och metadata (FÖRBÄTTRING 5)."""
    value: float
    decimals: Optional[int] = None  # None = INF (exakt)
    scale: int = 0
    unit: str = "SEK"
    sign: str = ""
    format: str = ""
    
    @property
    def precision_str(self) -> str:
        if self.decimals is None:
            return "exakt"
        elif self.decimals == 0:
            return "heltal"
        else:
            return f"{self.decimals} decimaler"

@dataclass
class Faststallelseintyg:
    """Strukturerat fastställelseintyg (FÖRBÄTTRING 6)."""
    intygsdatum: Optional[str] = None
    arsstamma_datum: Optional[str] = None
    resultatrakning_faststalld: bool = False
    balansrakning_faststalld: bool = False
    koncernresultatrakning_faststalld: Optional[bool] = None
    koncernbalansrakning_faststalld: Optional[bool] = None
    utdelning_per_aktie: Optional[float] = None
    utdelning_totalt: Optional[int] = None
    balanseras_i_ny_rakning: Optional[int] = None
    undertecknare: List[str] = field(default_factory=list)
    overensstammer_med_original: bool = False
    ar_esef: bool = False

@dataclass
class UtokadInformation:
    """Information taggad med utökad information-taxonomin (FÖRBÄTTRING 9)."""
    odefinierade_begrepp: List[Dict[str, Any]] = field(default_factory=list)
    andrade_rubriker: List[Dict[str, str]] = field(default_factory=list)
    notkopplingar: List[Dict[str, Any]] = field(default_factory=list)
    ar_fullstandigt_taggad: bool = True


# =============================================================================
# TAXONOMI-KONSTANTER (FÖRBÄTTRING 2)
# =============================================================================

TAXONOMI_VERSIONER = {
    # K2 Årsredovisning
    'k2/2024-09-12': TaxonomiInfo('2024-09-12', 'K2', 'se/fr/gaap/k2/2024-09-12', False),
    'k2/2021-10-31': TaxonomiInfo('2021-10-31', 'K2', 'se/fr/gaap/k2/2021-10-31', False),
    'k2/2017-09-30': TaxonomiInfo('2017-09-30', 'K2', 'se/fr/gaap/k2/2017-09-30', True, 
        'VARNING: K2 2017-09-30 är arkiverad och stöds EJ av Bolagsverket sedan 2023'),
    # K3 Årsredovisning
    'k3/2021-10-31': TaxonomiInfo('2021-10-31', 'K3', 'se/fr/gaap/k3/2021-10-31', False),
    'k3/2020-12-01': TaxonomiInfo('2020-12-01', 'K3', 'se/fr/gaap/k3/2020-12-01', False),
    'k3/2018-12-17': TaxonomiInfo('2018-12-17', 'K3', 'se/fr/gaap/k3/2018-12-17', True,
        'VARNING: K3 2018-12-17 är arkiverad och stöds EJ av Bolagsverket sedan 2023'),
    # K3K Koncernredovisning
    'k3k/2021-10-31': TaxonomiInfo('2021-10-31', 'K3K', 'se/fr/gaap/k3k/2021-10-31', False),
    'k3k/2020-12-01': TaxonomiInfo('2020-12-01', 'K3K', 'se/fr/gaap/k3k/2020-12-01', False),
    # Revisionsberättelse
    'ar/2020-12-01': TaxonomiInfo('2020-12-01', 'REVISION', 'se/fr/ar/2020-12-01', False),
    # Fastställelseintyg
    'ci/2022-09-01': TaxonomiInfo('2022-09-01', 'FASTSTALLELSE_ESEF', 'se/fr/ci/2022-09-01', False),
    'ci/2020-12-01': TaxonomiInfo('2020-12-01', 'FASTSTALLELSE', 'se/fr/ci/2020-12-01', False),
}


# =============================================================================
# K3K KONCERNBEGREPP (FÖRBÄTTRING 1)
# =============================================================================

K3K_BEGREPP = {
    'koncern_nettoomsattning': ['KoncernensNettoomsattning', 'NettoomsattningKoncern'],
    'koncern_rorelseresultat': ['KoncernensRorelseresultat', 'RorelseresultatKoncern'],
    'koncern_resultat_efter_finansiella': ['KoncernensResultatEfterFinansiellaPoster'],
    'koncern_arets_resultat': ['KoncernensAretsResultat', 'AretsResultatKoncern'],
    'koncern_goodwill': ['Goodwill', 'KoncernmassigGoodwill'],
    'koncern_eget_kapital': ['KoncernensEgetKapital', 'EgetKapitalKoncern'],
    'minoritetsandel': ['MinoritetsandelEgetKapital', 'Minoritetsintresse'],
    'koncern_balansomslutning': ['KoncernensSummaTillgangar', 'TillgangarKoncern'],
}


# =============================================================================
# BAS-KONTOMAPPNING (FÖRBÄTTRING 3)
# =============================================================================

BAS_KONTO_MAPPNING = {
    # Tillgångar
    'ImmateriellAnlaggningstillgangar': '10xx',
    'BalanseradeUtgifterUtvecklingsarbetenLiknande': '1010-1019',
    'Goodwill': '1050-1059',
    'MateriellaAnlaggningstillgangar': '11xx-12xx',
    'ByggnaderMark': '11xx',
    'MaskinerAndraTekniskaAnlaggningar': '12xx',
    'InventarierVerktygInstallationer': '1220-1229',
    'FinansiellaAnlaggningstillgangar': '13xx',
    'AndelarKoncernforetag': '1310-1319',
    'VarulagerMm': '14xx',
    'Kundfordringar': '1510-1519',
    'KassaBank': '19xx',
    # Eget kapital
    'Aktiekapital': '2081',
    'BalanseratResultat': '2091-2098',
    'AretsResultatEgetKapital': '2099',
    # Skulder
    'LangfristigaSkulder': '23xx-24xx',
    'KortfristigaSkulder': '25xx-29xx',
    'Leverantorsskulder': '2440-2449',
    # Resultaträkning - Intäkter
    'Nettoomsattning': '30xx-37xx',
    'OvrigaRorelseintakter': '39xx',
    # Resultaträkning - Kostnader
    'HandelsvarorKostnader': '40xx-49xx',
    'OvrigaExternaKostnader': '50xx-69xx',
    'Personalkostnader': '70xx-76xx',
    'AvskrivningarNedskrivningarMateriellaImmateriellaAnlaggningstillgangar': '78xx',
    # Finansiella poster
    'FinansiellaIntakter': '80xx-82xx',
    'FinansiellaKostnader': '83xx-84xx',
    'SkattAretsResultat': '89xx',
}

def get_bas_kontogrupp(begrepp: str) -> Optional[str]:
    """Hämta BAS-kontogrupp för ett taxonomi-begrepp."""
    if begrepp in BAS_KONTO_MAPPNING:
        return BAS_KONTO_MAPPNING[begrepp]
    for key, value in BAS_KONTO_MAPPNING.items():
        if key.lower() in begrepp.lower() or begrepp.lower() in key.lower():
            return value
    return None


# =============================================================================
# CACHE-HANTERING (FÖRBÄTTRING 8)
# =============================================================================

class CacheManager:
    """SQLite-baserad cache med TTL."""
    
    DEFAULT_TTL = {
        'arsredovisning': 30 * 24 * 3600,  # 30 dagar
        'company_info': 24 * 3600,          # 1 dag
        'dokumentlista': 7 * 24 * 3600,     # 7 dagar
        'ixbrl_document': 30 * 24 * 3600,   # 30 dagar
        'nyckeltal': 30 * 24 * 3600,        # 30 dagar
    }
    
    def __init__(self, db_path: str = None):
        if db_path is None:
            db_path = Path.home() / '.cache' / 'bolagsverket_mcp' / 'cache.db'
        self.db_path = Path(db_path)
        self.db_path.parent.mkdir(parents=True, exist_ok=True)
        self._init_db()
    
    def _init_db(self):
        with sqlite3.connect(self.db_path) as conn:
            conn.execute('''
                CREATE TABLE IF NOT EXISTS cache (
                    key TEXT PRIMARY KEY,
                    value TEXT,
                    category TEXT,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    expires_at TIMESTAMP,
                    hit_count INTEGER DEFAULT 0
                )
            ''')
            conn.execute('CREATE INDEX IF NOT EXISTS idx_expires ON cache(expires_at)')
    
    def get(self, category: str, identifier: str) -> Optional[Any]:
        key = f"{category}:{identifier}"
        with sqlite3.connect(self.db_path) as conn:
            cursor = conn.execute('''
                SELECT value FROM cache 
                WHERE key = ? AND expires_at > datetime('now')
            ''', (key,))
            row = cursor.fetchone()
            if row:
                conn.execute('UPDATE cache SET hit_count = hit_count + 1 WHERE key = ?', (key,))
                try:
                    return json.loads(row[0])
                except json.JSONDecodeError:
                    return row[0]
        return None
    
    def set(self, category: str, identifier: str, value: Any, ttl: int = None) -> None:
        if ttl is None:
            ttl = self.DEFAULT_TTL.get(category, 3600)
        key = f"{category}:{identifier}"
        expires_at = datetime.now() + timedelta(seconds=ttl)
        value_str = json.dumps(value, ensure_ascii=False, default=str) if isinstance(value, (dict, list)) else str(value)
        with sqlite3.connect(self.db_path) as conn:
            conn.execute('''
                INSERT OR REPLACE INTO cache (key, value, category, expires_at)
                VALUES (?, ?, ?, ?)
            ''', (key, value_str, category, expires_at))
    
    def delete(self, category: str, identifier: str) -> bool:
        key = f"{category}:{identifier}"
        with sqlite3.connect(self.db_path) as conn:
            cursor = conn.execute('DELETE FROM cache WHERE key = ?', (key,))
            return cursor.rowcount > 0
    
    def clear_expired(self) -> int:
        with sqlite3.connect(self.db_path) as conn:
            cursor = conn.execute("DELETE FROM cache WHERE expires_at <= datetime('now')")
            return cursor.rowcount
    
    def clear_all(self) -> int:
        with sqlite3.connect(self.db_path) as conn:
            cursor = conn.execute('DELETE FROM cache')
            return cursor.rowcount
    
    def get_stats(self) -> Dict[str, Any]:
        with sqlite3.connect(self.db_path) as conn:
            stats = {}
            cursor = conn.execute('SELECT COUNT(*) FROM cache')
            stats['total_entries'] = cursor.fetchone()[0]
            cursor = conn.execute('SELECT category, COUNT(*), SUM(hit_count) FROM cache GROUP BY category')
            stats['categories'] = {row[0]: {'count': row[1], 'hits': row[2] or 0} for row in cursor.fetchall()}
            cursor = conn.execute("SELECT COUNT(*) FROM cache WHERE expires_at <= datetime('now')")
            stats['expired_entries'] = cursor.fetchone()[0]
            stats['db_size_bytes'] = self.db_path.stat().st_size if self.db_path.exists() else 0
            return stats

# Global cache-instans
cache_manager = CacheManager()


# =============================================================================
# Token-hantering
# =============================================================================

class TokenManager:
    def __init__(self):
        self.access_token: Optional[str] = None
        self.token_expiry: Optional[datetime] = None
    
    def get_token(self, force_refresh: bool = False) -> str:
        if not force_refresh and self.access_token and self.token_expiry:
            if datetime.now() < self.token_expiry:
                return self.access_token
        
        logger.info("Hämtar ny OAuth2-token...")
        
        with httpx.Client(timeout=30.0) as client:
            response = client.post(
                TOKEN_URL,
                headers={"Content-Type": "application/x-www-form-urlencoded"},
                data={
                    "grant_type": "client_credentials",
                    "client_id": CLIENT_ID,
                    "client_secret": CLIENT_SECRET,
                    "scope": SCOPE
                }
            )
        
        if response.status_code != 200:
            logger.error(f"Token-fel: {response.status_code}")
            raise Exception(f"Token-fel: {response.status_code} - {response.text}")
        
        data = response.json()
        self.access_token = data["access_token"]
        expires_in = data.get("expires_in", 3600)
        self.token_expiry = datetime.now() + timedelta(seconds=expires_in - 60)
        
        logger.info("Token hämtad, giltig i %d sekunder", expires_in)
        return self.access_token


token_manager = TokenManager()


# =============================================================================
# FÖRBÄTTRING #7: ServerCapabilities
# =============================================================================

mcp = FastMCP("bolagsverket")


# =============================================================================
# API-hjälpfunktioner
# =============================================================================

def clean_org_nummer(org_nummer: str) -> str:
    return org_nummer.replace("-", "").replace(" ", "")


def format_org_nummer(org_nummer: str) -> str:
    clean = clean_org_nummer(org_nummer)
    if len(clean) == 10:
        return f"{clean[:6]}-{clean[6:]}"
    return clean


def luhn_checksum(number: str) -> bool:
    """Validera kontrollsiffra med Luhn-algoritmen.
    
    Bolagsverkets API kräver giltig kontrollsiffra och returnerar
    'Identitetsbeteckning har ogiltig kontrollsiffra.' vid fel.
    
    Algoritmen:
    1. Dubblera varannan siffra från höger (börja med näst sista)
    2. Om resultatet > 9, subtrahera 9
    3. Summera alla siffror
    4. Om summan är delbar med 10 är numret giltigt
    """
    digits = [int(d) for d in number]
    # Dubblera varannan siffra från höger (index -2, -4, -6, ...)
    for i in range(len(digits) - 2, -1, -2):
        digits[i] *= 2
        if digits[i] > 9:
            digits[i] -= 9
    return sum(digits) % 10 == 0


def validate_org_nummer(org_nummer: str) -> Tuple[bool, str]:
    """Validera organisationsnummer inklusive Luhn-kontroll.
    
    Organisationsnummer i Sverige:
    - 10 siffror för företag (NNNNNN-NNNN)
    - 12 siffror för personnummer (ÅÅÅÅMMDD-NNNN)
    
    Tredje siffran måste vara >= 2 för organisationsnummer
    (skiljer från personnummer där månad är 01-12).
    """
    clean = clean_org_nummer(org_nummer)
    
    if not clean.isdigit():
        return False, "Organisationsnummer får endast innehålla siffror"
    
    if len(clean) == 12:
        # Personnummer - ta bort sekelsiffror för validering
        # ÅÅÅÅMMDDNNNN -> ÅÅMMDDNNNN
        clean = clean[2:]
    
    if len(clean) != 10:
        return False, "Organisationsnummer måste vara 10 eller 12 siffror"
    
    # Luhn-validering
    if not luhn_checksum(clean):
        return False, "Organisationsnummer har ogiltig kontrollsiffra"
    
    return True, clean


def make_api_request(method: str, endpoint: str, json_body: Optional[Dict] = None) -> Dict[str, Any]:
    """Gör API-anrop till Bolagsverket.
    
    Hanterar autentisering och felformat enligt API-specifikationen.
    API:et returnerar ApiError-objekt med fälten: type, instance, timestamp,
    requestId, status, title, detail.
    """
    token = token_manager.get_token()
    
    request_id = str(uuid.uuid4())
    headers = {
        "Authorization": f"Bearer {token}",
        "X-Request-Id": request_id,
        "Accept": "application/json"
    }
    
    # POST-anrop kräver alltid Content-Type enligt Swagger-spec
    if method == "POST":
        headers["Content-Type"] = "application/json"
    
    url = f"{BASE_URL}{endpoint}"
    logger.debug(f"API-anrop: {method} {endpoint} (request_id: {request_id})")
    
    with httpx.Client(timeout=30.0) as client:
        if method == "GET":
            response = client.get(url, headers=headers)
        elif method == "POST":
            response = client.post(url, headers=headers, json=json_body)
        else:
            raise ValueError(f"Okänd HTTP-metod: {method}")
    
    if response.status_code != 200:
        # Hantera API:ets strukturerade felformat (ApiError enligt Swagger)
        # Fält: type, instance, timestamp, requestId, status, title, detail
        logger.error(f"API-fel: {response.status_code} - {response.text[:500]}")
        try:
            error_data = response.json()
            title = error_data.get("title", "Error")
            detail = error_data.get("detail", f"HTTP {response.status_code}")
            status = error_data.get("status", response.status_code)
            api_request_id = error_data.get("requestId", request_id)
            
            error_msg = f"{title}: {detail}"
            if status == 400:
                error_msg = f"Ogiltig begäran: {detail}"
            elif status == 401:
                error_msg = f"Ej autentiserad: {detail}"
            elif status == 403:
                error_msg = f"Åtkomst nekad: {detail}"
            elif status == 404:
                error_msg = f"Ej funnen: {detail}"
            elif status == 500:
                error_msg = f"Serverfel hos Bolagsverket: {detail}"
            
            logger.error(f"ApiError - requestId: {api_request_id}, status: {status}, title: {title}")
            raise Exception(error_msg)
        except json.JSONDecodeError:
            raise Exception(f"HTTP {response.status_code}: {response.text[:200]}")
    
    return response.json()


def download_document_bytes(dokument_id: str) -> bytes:
    """Ladda ner dokument från Bolagsverket.
    
    Returnerar ZIP-fil som bytes enligt API-specifikationen.
    """
    token = token_manager.get_token()
    
    request_id = str(uuid.uuid4())
    headers = {
        "Authorization": f"Bearer {token}",
        "Accept": "application/zip",  # Krävs enligt Swagger-spec
        "X-Request-Id": request_id
    }
    
    url = f"{BASE_URL}/dokument/{dokument_id}"
    logger.info(f"Laddar ner dokument: {dokument_id} (request_id: {request_id})")
    
    with httpx.Client(timeout=30.0) as client:  # 30 sek räcker för ZIP-filer
        response = client.get(url, headers=headers)
    
    if response.status_code != 200:
        # Hantera API:ets strukturerade felformat (ApiError)
        try:
            error_data = response.json()
            title = error_data.get("title", "Error")
            detail = error_data.get("detail", f"HTTP {response.status_code}")
            raise Exception(f"{title}: {detail}")
        except json.JSONDecodeError:
            raise Exception(f"HTTP {response.status_code}: Kunde inte ladda ner dokument")
    
    return response.content


def ensure_output_dir() -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    return OUTPUT_DIR


def fetch_company_info(org_nummer: str) -> CompanyInfo:
    """Hämta och strukturera företagsinformation.
    
    Förbättring #7: Extraherar alla fält från API-svaret inkl:
    - registreringsland
    - namnskyddslopnummer  
    - infört_hos_scb
    - alla_namn (alla registrerade namn)
    - datakalla_fel (fel från datakällor)
    
    Förbättring #8: Hanterar flera organisationer per identitetsbeteckning
    (vanligt för enskilda firmor med samma personnummer).
    """
    clean_nr = clean_org_nummer(org_nummer)
    data = make_api_request("POST", "/organisationer", {"identitetsbeteckning": clean_nr})
    
    orgs = data.get("organisationer", [])
    if not orgs:
        raise Exception(f"Företaget {org_nummer} hittades inte")
    
    # Förbättring #8: Logga om flera organisationer finns
    if len(orgs) > 1:
        logger.info(f"Identitetsbeteckning {clean_nr} har {len(orgs)} registrerade organisationer")
    
    org = orgs[0]
    
    # Förbättring #7: Samla datakällfel
    datakalla_fel = []
    
    def check_datakalla_fel(data_obj: dict, faltnamn: str) -> None:
        """Kontrollera om en datakälla returnerade fel."""
        if data_obj and data_obj.get("fel"):
            fel = data_obj["fel"]
            datakalla_fel.append({
                "falt": faltnamn,
                "typ": fel.get("typ", "OKÄNT"),
                "beskrivning": fel.get("felBeskrivning", "Okänt fel"),
                "dataproducent": data_obj.get("dataproducent", "Okänd")
            })
    
    # Kontrollera datakällfel för alla fält
    check_datakalla_fel(org.get("organisationsnamn", {}), "organisationsnamn")
    check_datakalla_fel(org.get("organisationsform", {}), "organisationsform")
    check_datakalla_fel(org.get("organisationsdatum", {}), "organisationsdatum")
    check_datakalla_fel(org.get("avregistreradOrganisation", {}), "avregistreradOrganisation")
    check_datakalla_fel(org.get("postadressOrganisation", {}), "postadressOrganisation")
    check_datakalla_fel(org.get("verksamhetsbeskrivning", {}), "verksamhetsbeskrivning")
    check_datakalla_fel(org.get("naringsgrenOrganisation", {}), "naringsgrenOrganisation")
    check_datakalla_fel(org.get("pagandeAvvecklingsEllerOmstruktureringsforfarande", {}), "pagandeForfarande")
    
    # Extrahera namn
    namn_data = org.get("organisationsnamn", {})
    namn_lista = namn_data.get("organisationsnamnLista", [])
    namn = namn_lista[0].get("namn", "Okänt") if namn_lista else "Okänt"
    
    # Förbättring #7: Alla namn (företagsnamn, särskilt företagsnamn, på främmande språk)
    alla_namn = []
    for n in namn_lista:
        alla_namn.append({
            "namn": n.get("namn"),
            "typ": n.get("organisationsnamntyp", {}).get("klartext"),
            "typ_kod": n.get("organisationsnamntyp", {}).get("kod"),
            "registreringsdatum": n.get("registreringsdatum"),
            "verksamhetsbeskrivning": n.get("verksamhetsbeskrivningSarskiltForetagsnamn")
        })
    
    avreg = org.get("avregistreradOrganisation", {})
    status = "Avregistrerad" if avreg and avreg.get("avregistreringsdatum") else "Aktiv"
    
    # Avregistreringsorsak
    avreg_orsak = org.get("avregistreringsorsak", {})
    avregistreringsorsak = avreg_orsak.get("klartext") if avreg_orsak else None
    
    adress_data = org.get("postadressOrganisation", {}).get("postadress", {}) or {}
    adress = {
        "utdelningsadress": adress_data.get("utdelningsadress", ""),
        "postnummer": adress_data.get("postnummer", ""),
        "postort": adress_data.get("postort", ""),
        "co_adress": adress_data.get("coAdress", ""),
        "land": adress_data.get("land", ""),  # Förbättring #7
    }
    
    sni_data = org.get("naringsgrenOrganisation", {}) or {}
    sni = sni_data.get("sni", []) or []
    sni_koder = [{"kod": s.get("kod", ""), "klartext": s.get("klartext", "")} for s in sni if s.get("kod")]
    
    # Pågående avveckling/omstrukturering (konkurs, likvidation)
    pagaende = org.get("pagandeAvvecklingsEllerOmstruktureringsforfarande", {})
    pagaende_lista = pagaende.get("pagandeAvvecklingsEllerOmstruktureringsforfarandeLista", []) if pagaende else []
    
    pagaende_konkurs = None
    pagaende_likvidation = None
    
    for p in pagaende_lista:
        kod = p.get("kod", "")
        if kod == "KK":  # Konkurs
            pagaende_konkurs = {
                "datum": p.get("fromDatum", ""),
                "typ": p.get("klartext", "Konkurs")
            }
        elif kod == "LI":  # Likvidation
            pagaende_likvidation = {
                "datum": p.get("fromDatum", ""),
                "typ": p.get("klartext", "Likvidation")
            }
    
    # Reklamspärr och verksam organisation
    reklamsparr_data = org.get("reklamsparr", {})
    reklamsparr = reklamsparr_data.get("kod") == "JA" if reklamsparr_data else None
    
    verksam_data = org.get("verksamOrganisation", {})
    verksam_organisation = verksam_data.get("kod") != "NEJ" if verksam_data else None
    
    # Förbättring #7: Nya fält
    registreringsland_data = org.get("registreringsland", {})
    registreringsland = None
    if registreringsland_data:
        registreringsland = {
            "kod": registreringsland_data.get("kod"),
            "klartext": registreringsland_data.get("klartext")
        }
    
    org_datum = org.get("organisationsdatum", {}) or {}
    infört_hos_scb = org_datum.get("infortHosScb")
    
    org_form = org.get("organisationsform", {}) or {}
    organisationsform_kod = org_form.get("kod")
    
    return CompanyInfo(
        org_nummer=format_org_nummer(clean_nr),
        namn=namn,
        organisationsform=org_form.get("klartext", "-"),
        juridisk_form=org.get("juridiskForm", {}).get("klartext") if org.get("juridiskForm") else None,
        registreringsdatum=org_datum.get("registreringsdatum", "-"),
        status=status,
        avregistreringsdatum=avreg.get("avregistreringsdatum") if avreg else None,
        avregistreringsorsak=avregistreringsorsak,
        adress=adress,
        verksamhet=org.get("verksamhetsbeskrivning", {}).get("beskrivning") if org.get("verksamhetsbeskrivning") else None,
        sni_koder=sni_koder,
        sate=org.get("sate", {}).get("lan") if org.get("sate") else None,
        pagaende_konkurs=pagaende_konkurs,
        pagaende_likvidation=pagaende_likvidation,
        reklamsparr=reklamsparr,
        verksam_organisation=verksam_organisation,
        # Nya fält (förbättring #7)
        registreringsland=registreringsland,
        namnskyddslopnummer=org.get("namnskyddslopnummer"),
        infört_hos_scb=infört_hos_scb,
        organisationsform_kod=organisationsform_kod,
        alla_namn=alla_namn if alla_namn else None,
        datakalla_fel=datakalla_fel if datakalla_fel else None,
    )


def fetch_all_companies_for_identity(org_nummer: str) -> List[CompanyInfo]:
    """Hämta ALLA organisationer för en identitetsbeteckning.
    
    Förbättring #8: För personnummer (enskilda firmor) kan samma person
    ha flera registrerade företag. Denna funktion returnerar alla.
    
    Returns:
        Lista med CompanyInfo för varje organisation
    """
    clean_nr = clean_org_nummer(org_nummer)
    data = make_api_request("POST", "/organisationer", {"identitetsbeteckning": clean_nr})
    
    orgs = data.get("organisationer", [])
    if not orgs:
        raise Exception(f"Inga organisationer hittades för {org_nummer}")
    
    companies = []
    for i, org in enumerate(orgs):
        # Skapa en CompanyInfo för varje organisation
        # Vi återanvänder logiken men måste extrahera per org
        try:
            # Temporärt sätt bara denna org i data
            temp_data = {"organisationer": [org]}
            # Använd intern extrahering
            company = _extract_company_info_from_org(org, clean_nr)
            companies.append(company)
        except Exception as e:
            logger.warning(f"Kunde inte extrahera organisation {i+1} för {clean_nr}: {e}")
    
    return companies


def _extract_company_info_from_org(org: Dict, clean_nr: str) -> CompanyInfo:
    """Extrahera CompanyInfo från ett organisation-objekt.
    
    Intern hjälpfunktion för att hantera flera organisationer.
    """
    # Samla datakällfel
    datakalla_fel = []
    
    def check_fel(data_obj, faltnamn):
        if data_obj and data_obj.get("fel"):
            fel = data_obj["fel"]
            datakalla_fel.append({
                "falt": faltnamn,
                "typ": fel.get("typ", "OKÄNT"),
                "beskrivning": fel.get("felBeskrivning", "Okänt fel")
            })
    
    # Extrahera namn
    namn_data = org.get("organisationsnamn", {})
    namn_lista = namn_data.get("organisationsnamnLista", [])
    namn = namn_lista[0].get("namn", "Okänt") if namn_lista else "Okänt"
    
    alla_namn = []
    for n in namn_lista:
        alla_namn.append({
            "namn": n.get("namn"),
            "typ": n.get("organisationsnamntyp", {}).get("klartext"),
            "typ_kod": n.get("organisationsnamntyp", {}).get("kod"),
            "registreringsdatum": n.get("registreringsdatum")
        })
    
    avreg = org.get("avregistreradOrganisation", {})
    status = "Avregistrerad" if avreg and avreg.get("avregistreringsdatum") else "Aktiv"
    
    avreg_orsak = org.get("avregistreringsorsak", {})
    avregistreringsorsak = avreg_orsak.get("klartext") if avreg_orsak else None
    
    adress_data = org.get("postadressOrganisation", {}).get("postadress", {}) or {}
    adress = {
        "utdelningsadress": adress_data.get("utdelningsadress", ""),
        "postnummer": adress_data.get("postnummer", ""),
        "postort": adress_data.get("postort", ""),
        "co_adress": adress_data.get("coAdress", ""),
        "land": adress_data.get("land", ""),
    }
    
    sni_data = org.get("naringsgrenOrganisation", {}) or {}
    sni = sni_data.get("sni", []) or []
    sni_koder = [{"kod": s.get("kod", ""), "klartext": s.get("klartext", "")} for s in sni if s.get("kod")]
    
    pagaende = org.get("pagandeAvvecklingsEllerOmstruktureringsforfarande", {})
    pagaende_lista = pagaende.get("pagandeAvvecklingsEllerOmstruktureringsforfarandeLista", []) if pagaende else []
    
    pagaende_konkurs = None
    pagaende_likvidation = None
    for p in pagaende_lista:
        kod = p.get("kod", "")
        if kod == "KK":
            pagaende_konkurs = {"datum": p.get("fromDatum", ""), "typ": p.get("klartext", "Konkurs")}
        elif kod == "LI":
            pagaende_likvidation = {"datum": p.get("fromDatum", ""), "typ": p.get("klartext", "Likvidation")}
    
    reklamsparr_data = org.get("reklamsparr", {})
    reklamsparr = reklamsparr_data.get("kod") == "JA" if reklamsparr_data else None
    
    verksam_data = org.get("verksamOrganisation", {})
    verksam_organisation = verksam_data.get("kod") != "NEJ" if verksam_data else None
    
    registreringsland_data = org.get("registreringsland", {})
    registreringsland = None
    if registreringsland_data:
        registreringsland = {
            "kod": registreringsland_data.get("kod"),
            "klartext": registreringsland_data.get("klartext")
        }
    
    org_datum = org.get("organisationsdatum", {}) or {}
    org_form = org.get("organisationsform", {}) or {}
    
    return CompanyInfo(
        org_nummer=format_org_nummer(clean_nr),
        namn=namn,
        organisationsform=org_form.get("klartext", "-"),
        juridisk_form=org.get("juridiskForm", {}).get("klartext") if org.get("juridiskForm") else None,
        registreringsdatum=org_datum.get("registreringsdatum", "-"),
        status=status,
        avregistreringsdatum=avreg.get("avregistreringsdatum") if avreg else None,
        avregistreringsorsak=avregistreringsorsak,
        adress=adress,
        verksamhet=org.get("verksamhetsbeskrivning", {}).get("beskrivning") if org.get("verksamhetsbeskrivning") else None,
        sni_koder=sni_koder,
        sate=org.get("sate", {}).get("lan") if org.get("sate") else None,
        pagaende_konkurs=pagaende_konkurs,
        pagaende_likvidation=pagaende_likvidation,
        reklamsparr=reklamsparr,
        verksam_organisation=verksam_organisation,
        registreringsland=registreringsland,
        namnskyddslopnummer=org.get("namnskyddslopnummer"),
        infört_hos_scb=org_datum.get("infortHosScb"),
        organisationsform_kod=org_form.get("kod"),
        alla_namn=alla_namn if alla_namn else None,
        datakalla_fel=datakalla_fel if datakalla_fel else None,
    )


# =============================================================================
# iXBRL Parser
# =============================================================================

class IXBRLParser:
    """Parser för iXBRL (Inline XBRL) årsredovisningar."""
    
    def __init__(self, xhtml_content: str):
        self.soup = BeautifulSoup(xhtml_content, 'lxml')
        self._cache = {}
    
    def _get_value(self, name_pattern: str, context: str = None, numeric: bool = True) -> Optional[Any]:
        """
        Extrahera värde från iXBRL-tagg.
        
        Hanterar:
        - ix:nonfraction för numeriska värden
        - ix:nonnumeric för text/datum
        - scale-attribut för skalade värden
        - sign-attribut för negativa värden (KRITISKT enligt tillämpningsanvisningar)
        - format-attribut för olika nummerformat
        
        Källa: taxonomier.se/tillampning.html, Exempel 1a-1e
        """
        tag_type = 'ix:nonfraction' if numeric else 'ix:nonnumeric'
        
        def name_match(x):
            return x and name_pattern.lower() in x.lower()
        
        attrs = {'name': name_match}
        if context:
            attrs['contextref'] = context
        
        tag = self.soup.find(tag_type, attrs)
        if tag:
            value = tag.text.strip()
            if numeric:
                # Hantera format-attribut enligt iXBRL-spec
                format_attr = tag.get('format', '')
                
                if 'numspacecomma' in format_attr or 'numspacedot' in format_attr:
                    # Mellanslag som tusentalsavgränsare
                    value = value.replace(' ', '')
                
                if 'commadecimal' in format_attr:
                    # Komma som decimaltecken (europeiskt format)
                    value = value.replace('.', '').replace(',', '.')
                elif 'dotdecimal' in format_attr:
                    # Punkt som decimaltecken (amerikanskt format)
                    value = value.replace(',', '')
                else:
                    # Standardhantering för svenskt format
                    value = value.replace(' ', '').replace(',', '.')
                
                # Hantera olika minus-tecken
                value = value.replace('−', '-').replace('–', '-')
                
                try:
                    scale = int(tag.get('scale', '0'))
                    numeric_value = float(value) * (10 ** scale)
                    
                    # KRITISKT: Hantera sign-attribut enligt tillämpningsanvisningar
                    # Källa: taxonomier.se/tillampning.html, Exempel 1d
                    sign = tag.get('sign', '')
                    if sign == '-':
                        numeric_value = -abs(numeric_value)
                    
                    return int(numeric_value)
                except ValueError:
                    return None
            return value
        return None
    
    def get_metadata(self) -> Dict[str, str]:
        return {
            'foretag_namn': self._get_value('ForetagetsNamn', numeric=False) or '',
            'org_nummer': self._get_value('Organisationsnummer', numeric=False) or '',
            'rakenskapsar_start': self._get_value('RakenskapsarForstaDag', numeric=False) or '',
            'rakenskapsar_slut': self._get_value('RakenskapsarSistaDag', numeric=False) or '',
            'undertecknat_datum': self._get_value('UndertecknandeDatum', numeric=False) or '',
            'sate': self._get_value('ForetagetsSate', numeric=False) or '',
        }
    
    def get_nyckeltal(self, period: str = 'period0') -> Nyckeltal:
        balans = period.replace('period', 'balans')
        
        nyckeltal = Nyckeltal(
            nettoomsattning=self._get_value('Nettoomsattning', period),
            resultat_efter_finansiella=self._get_value('ResultatEfterFinansiellaPoster', period),
            arets_resultat=self._get_value('AretsResultat', period),
            eget_kapital=self._get_value('EgetKapital', balans),
            balansomslutning=self._get_value('Tillgangar', balans) or self._get_value('SummaEgetKapitalSkulder', balans),
            soliditet=self._get_value('Soliditet', balans),
            antal_anstallda=self._get_value('MedelantalAnstallda', period),
        )
        
        nyckeltal.berakna_nyckeltal()
        return nyckeltal
    
    def get_flerarsoversikt(self) -> Dict[str, Nyckeltal]:
        oversikt = {}
        for i in range(4):
            period = f'period{i}'
            nyckeltal = self.get_nyckeltal(period)
            if nyckeltal.nettoomsattning is not None:
                oversikt[period] = nyckeltal
        return oversikt
    
    def get_personer(self) -> List[Person]:
        personer = []
        seen = set()
        
        patterns = [
            ('UnderskriftFaststallelseintygForetradareTilltalsnamn', 
             'UnderskriftFaststallelseintygForetradareEfternamn',
             'UnderskriftFaststallelseintygForetradareForetradarroll'),
            ('UnderskriftHandlingTilltalsnamn', 'UnderskriftHandlingEfternamn', None),
            ('UnderskriftRevisionsberattelseRevisorTilltalsnamn', 
             'UnderskriftRevisionsberattelseRevisorEfternamn',
             'UnderskriftRevisionsberattelseRevisorTitel'),
        ]
        
        for fornamn_pat, efternamn_pat, roll_pat in patterns:
            for tag in self.soup.find_all('ix:nonnumeric', {'name': lambda x: x and fornamn_pat in x}):
                fornamn = tag.text.strip()
                
                tuple_ref = tag.get('tupleref')
                efternamn = ''
                roll = ''
                
                if tuple_ref:
                    efternamn_tag = self.soup.find('ix:nonnumeric', {
                        'name': lambda x: x and efternamn_pat in x,
                        'tupleref': tuple_ref
                    })
                    if efternamn_tag:
                        efternamn = efternamn_tag.text.strip()
                    
                    if roll_pat:
                        roll_tag = self.soup.find('ix:nonnumeric', {
                            'name': lambda x: x and roll_pat in x,
                            'tupleref': tuple_ref
                        })
                        if roll_tag:
                            roll = roll_tag.text.strip()
                
                if not roll:
                    if 'Revisor' in fornamn_pat:
                        roll = 'Revisor'
                    elif 'Foretradar' in fornamn_pat:
                        roll = 'Företrädare'
                    else:
                        roll = 'Styrelseledamot'
                
                key = (fornamn, efternamn, roll)
                if key not in seen and fornamn:
                    seen.add(key)
                    personer.append(Person(fornamn=fornamn, efternamn=efternamn, roll=roll))
        
        return personer
    
    def get_balansrakning(self, period: str = 'balans0') -> Dict[str, Any]:
        return {
            'tillgangar': {
                'immateriella': self._get_value('ImmateriellAnlaggningstillgangar', period),
                'materiella': self._get_value('MateriellaAnlaggningstillgangar', period),
                'finansiella': self._get_value('FinansiellaAnlaggningstillgangar', period),
                'varulager': self._get_value('VarulagerMm', period),
                'kundfordringar': self._get_value('Kundfordringar', period),
                'kassa_bank': self._get_value('KassaBank', period),
                'summa_omsattning': self._get_value('Omsattningstillgangar', period),
                'summa_tillgangar': self._get_value('Tillgangar', period),
            },
            'eget_kapital_skulder': {
                'aktiekapital': self._get_value('Aktiekapital', period),
                'balanserat_resultat': self._get_value('BalanseratResultat', period),
                'arets_resultat': self._get_value('AretsResultatEgetKapital', period),
                'summa_eget_kapital': self._get_value('EgetKapital', period),
                'langfristiga_skulder': self._get_value('LangfristigaSkulder', period),
                'kortfristiga_skulder': self._get_value('KortfristigaSkulder', period),
                'leverantorsskulder': self._get_value('Leverantorsskulder', period),
                'summa_skulder': self._get_value('Skulder', period),
            }
        }
    
    def get_resultatrakning(self, period: str = 'period0') -> Dict[str, Any]:
        return {
            'nettoomsattning': self._get_value('Nettoomsattning', period),
            'ovriga_rorelseinktakter': self._get_value('OvrigaRorelseintakter', period),
            'summa_intakter': self._get_value('RorelseintakterLagerforandringarMm', period),
            'varor_handelsvaror': self._get_value('HandelsvarorKostnader', period),
            'ovriga_externa_kostnader': self._get_value('OvrigaExternaKostnader', period),
            'personalkostnader': self._get_value('Personalkostnader', period),
            'avskrivningar': self._get_value('AvskrivningarNedskrivningarMateriellaImmateriellaAnlaggningstillgangar', period),
            'rorelseresultat': self._get_value('Rorelseresultat', period),
            'finansiella_intakter': self._get_value('FinansiellaIntakter', period),
            'finansiella_kostnader': self._get_value('FinansiellaKostnader', period),
            'resultat_efter_finansiella': self._get_value('ResultatEfterFinansiellaPoster', period),
            'skatt': self._get_value('SkattAretsResultat', period),
            'arets_resultat': self._get_value('AretsResultat', period),
        }
    
    def get_personer_detaljerad(self) -> Tuple[List[Person], List[Person], Optional[Person]]:
        """Extrahera styrelse, revisorer och VD separat."""
        styrelse = []
        revisorer = []
        vd = None
        seen = set()
        
        # Styrelseledamöter
        for tag in self.soup.find_all('ix:nonnumeric', {'name': lambda x: x and 'UnderskriftFaststallelseintygForetradareTilltalsnamn' in x}):
            fornamn = tag.text.strip()
            tuple_ref = tag.get('tupleref')
            efternamn = ''
            roll = 'Styrelseledamot'
            
            if tuple_ref:
                efternamn_tag = self.soup.find('ix:nonnumeric', {
                    'name': lambda x: x and 'UnderskriftFaststallelseintygForetradareEfternamn' in x,
                    'tupleref': tuple_ref
                })
                if efternamn_tag:
                    efternamn = efternamn_tag.text.strip()
                
                roll_tag = self.soup.find('ix:nonnumeric', {
                    'name': lambda x: x and 'UnderskriftFaststallelseintygForetradareForetradarroll' in x,
                    'tupleref': tuple_ref
                })
                if roll_tag:
                    roll = roll_tag.text.strip()
            
            key = (fornamn, efternamn)
            if key not in seen and fornamn:
                seen.add(key)
                person = Person(fornamn=fornamn, efternamn=efternamn, roll=roll)
                
                if 'vd' in roll.lower() or 'verkställande' in roll.lower():
                    vd = person
                elif 'ordförande' in roll.lower():
                    styrelse.insert(0, person)
                else:
                    styrelse.append(person)
        
        # Revisorer
        for tag in self.soup.find_all('ix:nonnumeric', {'name': lambda x: x and 'UnderskriftRevisionsberattelseRevisorTilltalsnamn' in x}):
            fornamn = tag.text.strip()
            tuple_ref = tag.get('tupleref')
            efternamn = ''
            roll = 'Revisor'
            
            if tuple_ref:
                efternamn_tag = self.soup.find('ix:nonnumeric', {
                    'name': lambda x: x and 'UnderskriftRevisionsberattelseRevisorEfternamn' in x,
                    'tupleref': tuple_ref
                })
                if efternamn_tag:
                    efternamn = efternamn_tag.text.strip()
                
                roll_tag = self.soup.find('ix:nonnumeric', {
                    'name': lambda x: x and 'UnderskriftRevisionsberattelseRevisorTitel' in x,
                    'tupleref': tuple_ref
                })
                if roll_tag:
                    roll = roll_tag.text.strip()
            
            key = (fornamn, efternamn)
            if key not in seen and fornamn:
                seen.add(key)
                revisorer.append(Person(fornamn=fornamn, efternamn=efternamn, roll=roll))
        
        return styrelse, revisorer, vd
    
    def _extract_text_content(self, name_pattern: str) -> Optional[str]:
        """Extrahera textinnehåll."""
        for tag in self.soup.find_all('ix:nonnumeric', {'name': lambda x: x and name_pattern.lower() in x.lower()}):
            text = tag.get_text(separator='\n', strip=True)
            if len(text) > 50:
                return text
        return None
    
    def get_forvaltningsberattelse(self) -> Dict[str, str]:
        """Extrahera förvaltningsberättelse."""
        fb = {}
        patterns = {
            'allman_verksamhet': ['AllmanUppgiftVerksamhet', 'VerksamhetenArt'],
            'vasentliga_handelser': ['VasentligaHandelserRakenskapsaret', 'VasentligaHandelser'],
            'framtida_utveckling': ['ForvantadFramtidaUtveckling'],
            'resultatdisposition': ['ForslasResultatdisposition', 'Resultatdisposition'],
        }
        for key, name_patterns in patterns.items():
            for pattern in name_patterns:
                text = self._extract_text_content(pattern)
                if text:
                    fb[key] = text
                    break
        return fb
    
    def analyze_roda_flaggor(self) -> List[RodFlagga]:
        """Analysera och identifiera röda flaggor/varningar."""
        flaggor = []
        nyckeltal = self.get_nyckeltal('period0')
        nyckeltal_prev = self.get_nyckeltal('period1')
        
        # 1. Negativt eget kapital
        if nyckeltal.eget_kapital and nyckeltal.eget_kapital < 0:
            flaggor.append(RodFlagga(
                typ="negativt_eget_kapital",
                allvarlighet="kritisk",
                beskrivning=f"Negativt eget kapital: {nyckeltal.eget_kapital:,} SEK",
                varde=nyckeltal.eget_kapital,
                rekommendation="Bolaget kan behöva upprätta kontrollbalansräkning enligt ABL 25 kap."
            ))
        
        # 2. Eget kapital under hälften av aktiekapitalet
        aktiekapital = self._get_value('Aktiekapital', 'balans0')
        if aktiekapital and nyckeltal.eget_kapital and nyckeltal.eget_kapital < aktiekapital / 2:
            flaggor.append(RodFlagga(
                typ="lagt_eget_kapital",
                allvarlighet="kritisk",
                beskrivning=f"Eget kapital ({nyckeltal.eget_kapital:,}) understiger hälften av aktiekapitalet ({aktiekapital:,})",
                varde={'eget_kapital': nyckeltal.eget_kapital, 'aktiekapital': aktiekapital},
                rekommendation="Kontrollbalansräkning kan vara påkallad"
            ))
        
        # 3. Kraftigt fallande omsättning (>20%)
        if nyckeltal.nettoomsattning and nyckeltal_prev.nettoomsattning and nyckeltal_prev.nettoomsattning > 0:
            forandring = ((nyckeltal.nettoomsattning - nyckeltal_prev.nettoomsattning) / nyckeltal_prev.nettoomsattning) * 100
            if forandring < -20:
                flaggor.append(RodFlagga(
                    typ="fallande_omsattning",
                    allvarlighet="varning",
                    beskrivning=f"Omsättningen har minskat med {abs(forandring):.1f}% jämfört med föregående år",
                    varde={'nu': nyckeltal.nettoomsattning, 'foregaende': nyckeltal_prev.nettoomsattning},
                    rekommendation="Analysera orsakerna till omsättningsminskningen"
                ))
        
        # 4. Låg soliditet (<20%)
        if nyckeltal.soliditet is not None:
            if nyckeltal.soliditet < 0:
                flaggor.append(RodFlagga(
                    typ="negativ_soliditet",
                    allvarlighet="kritisk",
                    beskrivning=f"Negativ soliditet: {nyckeltal.soliditet:.1f}%",
                    varde=nyckeltal.soliditet,
                    rekommendation="Bolaget har mer skulder än tillgångar"
                ))
            elif nyckeltal.soliditet < 20:
                flaggor.append(RodFlagga(
                    typ="lag_soliditet",
                    allvarlighet="varning",
                    beskrivning=f"Låg soliditet: {nyckeltal.soliditet:.1f}%",
                    varde=nyckeltal.soliditet,
                    rekommendation="Överväg att stärka det egna kapitalet"
                ))
        
        # 5. Negativa resultat flera år i rad
        negativa_ar = 0
        for i in range(4):
            nt = self.get_nyckeltal(f'period{i}')
            if nt.arets_resultat is not None and nt.arets_resultat < 0:
                negativa_ar += 1
        
        if negativa_ar >= 2:
            flaggor.append(RodFlagga(
                typ="upprepade_forluster",
                allvarlighet="varning" if negativa_ar == 2 else "kritisk",
                beskrivning=f"Negativt resultat {negativa_ar} år i rad",
                varde=negativa_ar,
                rekommendation="Analysera lönsamheten och vidta åtgärder"
            ))
        
        # 6. Hög skuldsättningsgrad (>3x)
        if nyckeltal.eget_kapital and nyckeltal.balansomslutning and nyckeltal.eget_kapital > 0:
            skulder = nyckeltal.balansomslutning - nyckeltal.eget_kapital
            skuldsattningsgrad = skulder / nyckeltal.eget_kapital
            if skuldsattningsgrad > 3:
                flaggor.append(RodFlagga(
                    typ="hog_skuldsattning",
                    allvarlighet="varning",
                    beskrivning=f"Hög skuldsättningsgrad: {skuldsattningsgrad:.1f}x",
                    varde=skuldsattningsgrad,
                    rekommendation="Bolaget är högt belånat"
                ))
        
        # 7. Negativ vinstmarginal
        if nyckeltal.vinstmarginal is not None and nyckeltal.vinstmarginal < -10:
            flaggor.append(RodFlagga(
                typ="negativ_vinstmarginal",
                allvarlighet="varning",
                beskrivning=f"Kraftigt negativ vinstmarginal: {nyckeltal.vinstmarginal:.1f}%",
                varde=nyckeltal.vinstmarginal,
                rekommendation="Intäkterna täcker inte kostnaderna"
            ))
        
        return flaggor
    
    def parse_full(self) -> Arsredovisning:
        metadata = self.get_metadata()
        
        return Arsredovisning(
            org_nummer=metadata['org_nummer'],
            foretag_namn=metadata['foretag_namn'],
            rakenskapsar_start=metadata['rakenskapsar_start'],
            rakenskapsar_slut=metadata['rakenskapsar_slut'],
            nyckeltal=self.get_nyckeltal(),
            personer=self.get_personer(),
            balansrakning=self.get_balansrakning(),
            resultatrakning=self.get_resultatrakning(),
            noter={},
            metadata=metadata,
        )
    
    def parse_full_extended(self) -> FullArsredovisning:
        """Parsa hela årsredovisningen till FullArsredovisning med alla detaljer."""
        metadata = self.get_metadata()
        styrelse, revisorer, vd = self.get_personer_detaljerad()
        flerarsoversikt = self.get_flerarsoversikt()
        
        return FullArsredovisning(
            org_nummer=metadata['org_nummer'],
            foretag_namn=metadata['foretag_namn'],
            rakenskapsar_start=metadata['rakenskapsar_start'],
            rakenskapsar_slut=metadata['rakenskapsar_slut'],
            nyckeltal=self.get_nyckeltal(),
            styrelse=styrelse,
            revisorer=revisorer,
            vd=vd,
            balansrakning=self.get_balansrakning(),
            resultatrakning=self.get_resultatrakning(),
            forvaltningsberattelse=self.get_forvaltningsberattelse(),
            noter={},
            flerarsdata=[asdict(nt) for nt in flerarsoversikt.values()],
            roda_flaggor=self.analyze_roda_flaggor(),
            metadata=metadata,
        )
    
    # =========================================================================
    # NYA METODER FÖR v5.1.0
    # =========================================================================
    
    def get_taxonomi_info(self) -> List[TaxonomiInfo]:
        """
        FÖRBÄTTRING 2: Identifiera vilka taxonomier dokumentet använder.
        
        Detekterar taxonomiversion från schemaRef och namespaces.
        Varnar för arkiverade versioner som ej stöds av Bolagsverket.
        """
        taxonomier = []
        doc_text = str(self.soup).lower()
        
        for tax_id, tax_info in TAXONOMI_VERSIONER.items():
            # Sök efter taxonomi-ID i dokumentet
            if tax_id.replace('/', '-') in doc_text or tax_id in doc_text:
                taxonomier.append(tax_info)
            # Sök efter entry point
            elif tax_info.entry_point.lower() in doc_text:
                taxonomier.append(tax_info)
        
        return taxonomier
    
    def get_taxonomi_varningar(self) -> List[str]:
        """Hämta varningar för använda taxonomier."""
        return [t.varning for t in self.get_taxonomi_info() if t.ar_arkiverad and t.varning]
    
    def ar_koncernredovisning(self) -> bool:
        """
        FÖRBÄTTRING 1: Kontrollera om dokumentet är en koncernredovisning.
        """
        # Kolla taxonomi-typ
        for tax in self.get_taxonomi_info():
            if tax.typ == 'K3K':
                return True
        
        # Sök efter koncernbegrepp
        for begrepp_lista in K3K_BEGREPP.values():
            for begrepp in begrepp_lista:
                if self._get_value(begrepp, numeric=True) is not None:
                    return True
        
        return False
    
    def get_koncern_nyckeltal(self, period: str = 'period0') -> Optional[KoncernNyckeltal]:
        """
        FÖRBÄTTRING 1: Extrahera koncernnyckeltal (endast för K3K).
        """
        if not self.ar_koncernredovisning():
            return None
        
        balans = period.replace('period', 'balans')
        nyckeltal = KoncernNyckeltal()
        
        for attr, begrepp_lista in K3K_BEGREPP.items():
            ctx = balans if 'kapital' in attr or 'balansomslutning' in attr else period
            for begrepp in begrepp_lista:
                value = self._get_value(begrepp, ctx, numeric=True)
                if value is not None:
                    # Mappa K3K-attribut till KoncernNyckeltal-attribut
                    mapped_attr = attr.replace('koncern_', '') if attr.startswith('koncern_') else attr
                    if hasattr(nyckeltal, mapped_attr):
                        setattr(nyckeltal, mapped_attr, value)
                    elif hasattr(nyckeltal, 'koncern_' + mapped_attr):
                        setattr(nyckeltal, 'koncern_' + mapped_attr, value)
                    break
        
        nyckeltal.berakna_koncern_nyckeltal()
        return nyckeltal
    
    def get_bas_mappning(self) -> Dict[str, Dict[str, Any]]:
        """
        FÖRBÄTTRING 3: Generera BAS-kontomappning för alla extraherade värden.
        """
        mappings = {}
        
        for tag in self.soup.find_all('ix:nonfraction'):
            name = tag.get('name', '')
            if ':' in name:
                name = name.split(':')[1]
            
            value = self._parse_numeric_tag(tag)
            if value is not None:
                bas = get_bas_kontogrupp(name)
                if bas:
                    mappings[name] = {
                        'varde': value,
                        'bas_konto': bas,
                        'context': tag.get('contextref', '')
                    }
        
        return mappings
    
    def _parse_numeric_tag(self, tag) -> Optional[int]:
        """Hjälpmetod för att parsa numeriskt värde."""
        if not tag:
            return None
        value = tag.text.strip().replace(' ', '').replace(',', '.').replace('−', '-').replace('–', '-')
        try:
            scale = int(tag.get('scale', '0'))
            numeric_value = float(value) * (10 ** scale)
            sign = tag.get('sign', '')
            if sign == '-':
                numeric_value = -abs(numeric_value)
            return int(numeric_value)
        except ValueError:
            return None
    
    def get_revisionsberattelse(self) -> Optional[Revisionsberattelse]:
        """
        FÖRBÄTTRING 4: Extrahera strukturerad revisionsberättelse.
        """
        rb = Revisionsberattelse()
        
        # Revisor-info
        rb.revisor_namn = self._get_value('UnderskriftRevisionsberattelseRevisorTilltalsnamn', numeric=False)
        if rb.revisor_namn:
            efternamn = self._get_value('UnderskriftRevisionsberattelseRevisorEfternamn', numeric=False)
            if efternamn:
                rb.revisor_namn = f"{rb.revisor_namn} {efternamn}"
        
        rb.revisor_titel = self._get_value('UnderskriftRevisionsberattelseRevisorTitel', numeric=False)
        rb.revisionsbolag = self._get_value('UnderskriftRevisionsberattelseRevisionsbolag', numeric=False)
        
        # Uttalanden
        rb.uttalande_arsredovisning = self._extract_text_content('RevisorsUttalandeOmArsredovisningen')
        rb.uttalande_forvaltning = self._extract_text_content('RevisorsUttalandeOmForvaltningen')
        rb.grund_for_uttalande = self._extract_text_content('GrundForUttalanden')
        
        # Anmärkningar
        anm = self._extract_text_content('AnmarkningarRevisionsberattelse')
        if anm:
            rb.anmarkningar.append(anm)
            rb.ar_ren = False
        
        # Datum och ort
        rb.datum = self._get_value('UnderskriftRevisionsberattelseDatum', numeric=False)
        rb.ort = self._get_value('UnderskriftRevisionsberattelseOrt', numeric=False)
        
        # Koncern
        if self.ar_koncernredovisning():
            rb.typ = 'koncern'
            rb.uttalande_koncernredovisning = self._extract_text_content('RevisorsUttalandeOmKoncernredovisningen')
        
        return rb if rb.revisor_namn or rb.datum else None
    
    def get_value_with_precision(self, name_pattern: str, context: str = None) -> Optional[PrecisionValue]:
        """
        FÖRBÄTTRING 5: Extrahera värde med full precision och metadata.
        """
        def name_match(x):
            return x and name_pattern.lower() in x.lower()
        
        attrs = {'name': name_match}
        if context:
            attrs['contextref'] = context
        
        tag = self.soup.find('ix:nonfraction', attrs)
        if not tag:
            return None
        
        value_text = tag.text.strip().replace(' ', '').replace(',', '.').replace('−', '-').replace('–', '-')
        
        decimals_attr = tag.get('decimals')
        decimals = None
        if decimals_attr and decimals_attr.upper() != 'INF':
            try:
                decimals = int(decimals_attr)
            except ValueError:
                pass
        
        try:
            scale = int(tag.get('scale', '0'))
            numeric_value = float(value_text) * (10 ** scale)
            sign = tag.get('sign', '')
            if sign == '-':
                numeric_value = -abs(numeric_value)
            
            return PrecisionValue(
                value=numeric_value,
                decimals=decimals,
                scale=scale,
                unit=tag.get('unitref', 'SEK'),
                sign=sign,
                format=tag.get('format', '')
            )
        except ValueError:
            return None
    
    def get_faststallelseintyg(self) -> Optional[Faststallelseintyg]:
        """
        FÖRBÄTTRING 6: Extrahera strukturerat fastställelseintyg.
        """
        fi = Faststallelseintyg()
        
        # Datum
        fi.intygsdatum = self._get_value('UnderskriftFaststallelseintygDatum', numeric=False)
        fi.arsstamma_datum = self._get_value('ArsstammaDatum', numeric=False) or \
                             self._get_value('FastallelseintygDatum', numeric=False)
        
        # Resultatdisposition
        fi.utdelning_per_aktie = self._get_value('UtdelningPerAktie', numeric=True)
        fi.utdelning_totalt = self._get_value('Utdelning', numeric=True)
        fi.balanseras_i_ny_rakning = self._get_value('BalanserasINyRakning', numeric=True)
        
        # Undertecknare
        for tag in self.soup.find_all('ix:nonnumeric', {
            'name': lambda x: x and 'UnderskriftFaststallelseintygForetradareTilltalsnamn' in x
        }):
            fornamn = tag.text.strip()
            tuple_ref = tag.get('tupleref')
            if tuple_ref:
                efternamn_tag = self.soup.find('ix:nonnumeric', {
                    'name': lambda x: x and 'UnderskriftFaststallelseintygForetradareEfternamn' in x,
                    'tupleref': tuple_ref
                })
                if efternamn_tag:
                    fi.undertecknare.append(f"{fornamn} {efternamn_tag.text.strip()}")
                else:
                    fi.undertecknare.append(fornamn)
        
        # Flaggor
        fi.resultatrakning_faststalld = bool(fi.arsstamma_datum)
        fi.balansrakning_faststalld = bool(fi.arsstamma_datum)
        
        if self.ar_koncernredovisning():
            fi.koncernresultatrakning_faststalld = fi.resultatrakning_faststalld
            fi.koncernbalansrakning_faststalld = fi.balansrakning_faststalld
        
        return fi if fi.intygsdatum or fi.undertecknare else None
    
    def get_utokad_information(self) -> UtokadInformation:
        """
        FÖRBÄTTRING 9: Extrahera information taggad med utökad information-taxonomin.
        """
        ui = UtokadInformation()
        
        # Odefinierade begrepp
        for tag in self.soup.find_all('ix:nonfraction', {
            'name': lambda x: x and 'OdefinieratBegrepp' in x
        }):
            ui.odefinierade_begrepp.append({
                'namn': tag.get('name', ''),
                'varde': self._parse_numeric_tag(tag),
                'context': tag.get('contextref', ''),
            })
        
        # Ändrade rubriker
        for tag in self.soup.find_all('ix:nonnumeric', {
            'name': lambda x: x and 'AndradRubrik' in x
        }):
            tuple_ref = tag.get('tupleref')
            if tuple_ref:
                ursprunglig = self.soup.find('ix:nonnumeric', {
                    'name': lambda x: x and 'UrsprungligRubrik' in x,
                    'tupleref': tuple_ref
                })
                ny = self.soup.find('ix:nonnumeric', {
                    'name': lambda x: x and 'NyRubrik' in x,
                    'tupleref': tuple_ref
                })
                if ursprunglig and ny:
                    ui.andrade_rubriker.append({
                        'ursprunglig': ursprunglig.text.strip(),
                        'ny': ny.text.strip(),
                    })
        
        # Notkopplingar
        for tag in self.soup.find_all('ix:nonnumeric', {
            'name': lambda x: x and 'Notkoppling' in x
        }):
            ui.notkopplingar.append({
                'begrepp': tag.get('name', ''),
                'not_nummer': tag.text.strip(),
                'context': tag.get('contextref', ''),
            })
        
        # Kontrollera flaggor
        flagga = self._get_value('ArsredovisningInnehallNotTaggadInformation', numeric=False)
        if flagga and flagga.lower() == 'true':
            ui.ar_fullstandigt_taggad = False
        
        return ui


def fetch_and_parse_arsredovisning(org_nummer: str, index: int = 0) -> Tuple[Arsredovisning, bytes, bytes]:
    """Hämta och parsa årsredovisning.
    
    Returnerar:
        Tuple med (Arsredovisning, xhtml_bytes, zip_bytes)
    """
    clean_nr = clean_org_nummer(org_nummer)
    
    dok_data = make_api_request("POST", "/dokumentlista", {"identitetsbeteckning": clean_nr})
    dokument = dok_data.get("dokument", [])
    
    if not dokument:
        raise Exception("Inga årsredovisningar hittades")
    
    if index >= len(dokument):
        raise Exception(f"Index {index} finns inte. Det finns {len(dokument)} årsredovisningar.")
    
    dok = dokument[index]
    dok_id = dok.get("dokumentId")
    
    logger.info(f"Hämtar årsredovisning {index+1}/{len(dokument)} för {format_org_nummer(clean_nr)}")
    
    zip_bytes = download_document_bytes(dok_id)
    
    xhtml_content = None
    xhtml_filename = None
    with zipfile.ZipFile(BytesIO(zip_bytes)) as zf:
        for name in zf.namelist():
            if name.lower().endswith(('.xhtml', '.html', '.xml')):
                xhtml_content = zf.read(name).decode('utf-8')
                xhtml_filename = name
                break
    
    if not xhtml_content:
        raise Exception("Ingen XHTML-fil hittades i ZIP-arkivet")
    
    parser = IXBRLParser(xhtml_content)
    return parser.parse_full(), xhtml_content.encode('utf-8'), zip_bytes


def fetch_full_arsredovisning(org_nummer: str, index: int = 0) -> Tuple[FullArsredovisning, bytes, bytes]:
    """Hämta och parsa KOMPLETT årsredovisning med alla detaljer.
    
    Returnerar:
        Tuple med (FullArsredovisning, xhtml_bytes, zip_bytes)
    """
    clean_nr = clean_org_nummer(org_nummer)
    
    dok_data = make_api_request("POST", "/dokumentlista", {"identitetsbeteckning": clean_nr})
    dokument = dok_data.get("dokument", [])
    
    if not dokument:
        raise Exception("Inga årsredovisningar hittades")
    
    if index >= len(dokument):
        raise Exception(f"Index {index} finns inte. Det finns {len(dokument)} årsredovisningar.")
    
    dok = dokument[index]
    dok_id = dok.get("dokumentId")
    
    logger.info(f"Hämtar KOMPLETT årsredovisning {index+1}/{len(dokument)} för {format_org_nummer(clean_nr)}")
    
    zip_bytes = download_document_bytes(dok_id)
    
    xhtml_content = None
    with zipfile.ZipFile(BytesIO(zip_bytes)) as zf:
        for name in zf.namelist():
            if name.lower().endswith(('.xhtml', '.html', '.xml')):
                xhtml_content = zf.read(name).decode('utf-8')
                break
    
    if not xhtml_content:
        raise Exception("Ingen XHTML-fil hittades i ZIP-arkivet")
    
    parser = IXBRLParser(xhtml_content)
    return parser.parse_full_extended(), xhtml_content.encode('utf-8'), zip_bytes

def export_to_json(data: Any) -> str:
    if hasattr(data, '__dataclass_fields__'):
        data = asdict(data)
    return json.dumps(data, indent=2, ensure_ascii=False, default=str)


def export_to_csv(data: Dict[str, Any], filename: str = None) -> str:
    output = StringIO()
    writer = csv.writer(output, delimiter=';')
    
    writer.writerow(['Nyckeltal', 'Värde', 'Enhet'])
    
    labels = {
        'nettoomsattning': ('Nettoomsättning', 'SEK'),
        'resultat_efter_finansiella': ('Resultat efter finansiella poster', 'SEK'),
        'arets_resultat': ('Årets resultat', 'SEK'),
        'eget_kapital': ('Eget kapital', 'SEK'),
        'balansomslutning': ('Balansomslutning', 'SEK'),
        'soliditet': ('Soliditet', '%'),
        'vinstmarginal': ('Vinstmarginal', '%'),
        'roe': ('Avkastning på eget kapital (ROE)', '%'),
        'antal_anstallda': ('Antal anställda', 'st'),
    }
    
    for key, (label, unit) in labels.items():
        value = data.get(key)
        if value is not None:
            writer.writerow([label, value, unit])
    
    csv_content = output.getvalue()
    
    if filename:
        filepath = ensure_output_dir() / filename
        with open(filepath, 'w', encoding='utf-8-sig') as f:
            f.write(csv_content)
        return str(filepath)
    
    return csv_content


def export_to_excel(arsredovisning: Arsredovisning, filename: str = None) -> str:
    if not EXCEL_AVAILABLE:
        return handle_error(ErrorCode.EXPORT_ERROR, "Excel-export ej tillgänglig", reason="openpyxl saknas")
    
    wb = openpyxl.Workbook()
    
    header_font = Font(bold=True, size=12)
    title_font = Font(bold=True, size=14)
    money_format = '#,##0'
    
    ws = wb.active
    ws.title = "Översikt"
    
    ws['A1'] = arsredovisning.foretag_namn
    ws['A1'].font = title_font
    ws['A2'] = f"Org.nr: {format_org_nummer(arsredovisning.org_nummer)}"
    ws['A3'] = f"Räkenskapsår: {arsredovisning.rakenskapsar_start} - {arsredovisning.rakenskapsar_slut}"
    
    ws['A5'] = 'Nyckeltal'
    ws['A5'].font = header_font
    
    row = 6
    nyckeltal = asdict(arsredovisning.nyckeltal)
    labels = {
        'nettoomsattning': 'Nettoomsättning',
        'resultat_efter_finansiella': 'Resultat efter finansiella poster',
        'arets_resultat': 'Årets resultat',
        'eget_kapital': 'Eget kapital',
        'balansomslutning': 'Balansomslutning',
        'soliditet': 'Soliditet (%)',
        'vinstmarginal': 'Vinstmarginal (%)',
        'roe': 'ROE (%)',
        'antal_anstallda': 'Antal anställda',
    }
    
    for key, label in labels.items():
        value = nyckeltal.get(key)
        if value is not None:
            ws[f'A{row}'] = label
            ws[f'B{row}'] = value
            if key not in ('soliditet', 'vinstmarginal', 'roe', 'antal_anstallda'):
                ws[f'B{row}'].number_format = money_format
            row += 1
    
    ws2 = wb.create_sheet("Personer")
    ws2['A1'] = 'Förnamn'
    ws2['B1'] = 'Efternamn'
    ws2['C1'] = 'Roll'
    for cell in ws2[1]:
        cell.font = header_font
    
    for i, person in enumerate(arsredovisning.personer, 2):
        ws2[f'A{i}'] = person.fornamn
        ws2[f'B{i}'] = person.efternamn
        ws2[f'C{i}'] = person.roll
    
    if not filename:
        clean_name = re.sub(r'[^\w\s-]', '', arsredovisning.foretag_namn)
        filename = f"{clean_name}_{arsredovisning.rakenskapsar_slut[:4]}.xlsx"
    
    filepath = ensure_output_dir() / filename
    wb.save(filepath)
    logger.info(f"Excel exporterad till: {filepath}")
    return str(filepath)


def export_to_pdf(arsredovisning: Arsredovisning, filename: str = None) -> str:
    if not PDF_AVAILABLE:
        return handle_error(ErrorCode.EXPORT_ERROR, "PDF-export ej tillgänglig", reason="weasyprint saknas")
    
    nyckeltal = arsredovisning.nyckeltal
    
    personer_html = ""
    for p in arsredovisning.personer:
        personer_html += f"<tr><td>{p.fornamn}</td><td>{p.efternamn}</td><td>{p.roll}</td></tr>"
    
    def fmt(val):
        return f"{val:,}" if val else "-"
    
    html_content = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="utf-8">
        <style>
            body {{ font-family: 'Helvetica Neue', Arial, sans-serif; margin: 40px; color: #333; }}
            h1 {{ color: #1a365d; border-bottom: 2px solid #1a365d; padding-bottom: 10px; }}
            h2 {{ color: #2c5282; margin-top: 30px; }}
            .info {{ background: #f7fafc; padding: 15px; border-radius: 5px; margin: 20px 0; }}
            table {{ width: 100%; border-collapse: collapse; margin: 20px 0; }}
            th, td {{ padding: 10px; text-align: left; border-bottom: 1px solid #e2e8f0; }}
            th {{ background: #edf2f7; font-weight: bold; }}
            .number {{ text-align: right; font-family: monospace; }}
            .highlight {{ background: #ebf8ff; }}
        </style>
    </head>
    <body>
        <h1>{arsredovisning.foretag_namn}</h1>
        <div class="info">
            <strong>Organisationsnummer:</strong> {format_org_nummer(arsredovisning.org_nummer)}<br>
            <strong>Räkenskapsår:</strong> {arsredovisning.rakenskapsar_start} – {arsredovisning.rakenskapsar_slut}
        </div>
        
        <h2>Nyckeltal</h2>
        <table>
            <tr><th>Nyckeltal</th><th class="number">Belopp (SEK)</th></tr>
            <tr><td>Nettoomsättning</td><td class="number">{fmt(nyckeltal.nettoomsattning)}</td></tr>
            <tr><td>Resultat efter finansiella poster</td><td class="number">{fmt(nyckeltal.resultat_efter_finansiella)}</td></tr>
            <tr class="highlight"><td><strong>Årets resultat</strong></td><td class="number"><strong>{fmt(nyckeltal.arets_resultat)}</strong></td></tr>
            <tr><td>Eget kapital</td><td class="number">{fmt(nyckeltal.eget_kapital)}</td></tr>
            <tr><td>Soliditet</td><td class="number">{nyckeltal.soliditet or '-'} %</td></tr>
            <tr><td>Vinstmarginal</td><td class="number">{nyckeltal.vinstmarginal or '-'} %</td></tr>
            <tr><td>ROE</td><td class="number">{nyckeltal.roe or '-'} %</td></tr>
        </table>
        
        <h2>Personer</h2>
        <table>
            <tr><th>Förnamn</th><th>Efternamn</th><th>Roll</th></tr>
            {personer_html}
        </table>
        
        <div style="margin-top: 40px; font-size: 11px; color: #718096;">
            Genererad: {datetime.now().strftime('%Y-%m-%d %H:%M')} | Källa: Bolagsverket
        </div>
    </body>
    </html>
    """
    
    if not filename:
        clean_name = re.sub(r'[^\w\s-]', '', arsredovisning.foretag_namn)
        filename = f"{clean_name}_{arsredovisning.rakenskapsar_slut[:4]}.pdf"
    
    filepath = ensure_output_dir() / filename
    HTML(string=html_content).write_pdf(filepath)
    logger.info(f"PDF exporterad till: {filepath}")
    return str(filepath)


def export_to_markdown(arsredovisning: Arsredovisning) -> str:
    nyckeltal = arsredovisning.nyckeltal
    
    lines = [
        f"# {arsredovisning.foretag_namn}",
        f"",
        f"**Organisationsnummer:** {format_org_nummer(arsredovisning.org_nummer)}  ",
        f"**Räkenskapsår:** {arsredovisning.rakenskapsar_start} – {arsredovisning.rakenskapsar_slut}",
        f"",
        f"## Nyckeltal",
        f"",
        f"| Nyckeltal | Belopp |",
        f"|-----------|--------|",
    ]
    
    if nyckeltal.nettoomsattning:
        lines.append(f"| Nettoomsättning | {nyckeltal.nettoomsattning:,} SEK |")
    if nyckeltal.resultat_efter_finansiella:
        lines.append(f"| Resultat efter finansiella poster | {nyckeltal.resultat_efter_finansiella:,} SEK |")
    if nyckeltal.arets_resultat:
        lines.append(f"| **Årets resultat** | **{nyckeltal.arets_resultat:,} SEK** |")
    if nyckeltal.eget_kapital:
        lines.append(f"| Eget kapital | {nyckeltal.eget_kapital:,} SEK |")
    if nyckeltal.soliditet:
        lines.append(f"| Soliditet | {nyckeltal.soliditet} % |")
    if nyckeltal.vinstmarginal:
        lines.append(f"| Vinstmarginal | {nyckeltal.vinstmarginal} % |")
    if nyckeltal.roe:
        lines.append(f"| ROE | {nyckeltal.roe} % |")
    
    lines.extend([
        f"",
        f"## Personer",
        f"",
        f"| Namn | Roll |",
        f"|------|------|",
    ])
    
    for p in arsredovisning.personer:
        lines.append(f"| {p.fullnamn} | {p.roll} |")
    
    return "\n".join(lines)


# =============================================================================
# FÖRBÄTTRING #6: Pydantic-modeller med förbättrade inputSchema
# =============================================================================

class ResponseFormat(str, Enum):
    MARKDOWN = "markdown"
    JSON = "json"
    CSV = "csv"
    EXCEL = "excel"
    PDF = "pdf"
    XHTML = "xhtml"  # Original iXBRL-fil
    ZIP = "zip"      # Original ZIP från Bolagsverket


# =============================================================================
# FÖRBÄTTRING v5.0: Mixin med field_validator för organisationsnummer
# =============================================================================

class OrgNummerMixin(BaseModel):
    """
    Mixin-klass med robust organisationsnummer-validering.
    
    Funktionalitet:
    - Accepterar format: "5566778899", "556677-8899", "556677 8899"
    - Normaliserar till 10 siffror utan bindestreck
    - Validerar med Luhn-algoritmen (kontrollsiffra)
    - Ger tydliga felmeddelanden på svenska
    """
    
    org_nummer: str = Field(
        min_length=10,
        max_length=13,
        description="10-siffrigt organisationsnummer (t.ex. 5567671267 eller 556767-1267)",
        json_schema_extra={"examples": ["5567671267", "556767-1267"]}
    )
    
    @field_validator('org_nummer', mode='before')
    @classmethod
    def validate_and_normalize_org_nummer(cls, v: str) -> str:
        """
        Validerar och normaliserar organisationsnummer.
        
        Steg:
        1. Ta bort bindestreck och mellanslag
        2. Kontrollera att det är 10 eller 12 siffror
        3. Validera kontrollsiffra med Luhn-algoritmen
        4. Returnera normaliserat 10-siffrigt nummer
        
        Raises:
            ValueError: Om organisationsnumret är ogiltigt
        """
        if not isinstance(v, str):
            raise ValueError(f"Organisationsnummer måste vara en sträng, fick: {type(v).__name__}")
        
        # Steg 1: Normalisera - ta bort bindestreck och mellanslag
        clean = re.sub(r'[-\s]', '', v.strip())
        
        # Steg 2: Kontrollera längd
        if not clean.isdigit():
            raise ValueError(
                f"Organisationsnummer får endast innehålla siffror, "
                f"bindestreck och mellanslag. Fick: '{v}'"
            )
        
        # Hantera personnummer (12 siffror) -> ta bort sekelsiffror
        if len(clean) == 12:
            clean = clean[2:]  # ÅÅÅÅMMDDNNNN -> ÅÅMMDDNNNN
        
        if len(clean) != 10:
            raise ValueError(
                f"Organisationsnummer måste vara 10 siffror (eller 12 för personnummer). "
                f"'{v}' har {len(clean)} siffror efter normalisering."
            )
        
        # Steg 3: Luhn-validering
        def luhn_check(number: str) -> bool:
            """Validera kontrollsiffra med Luhn-algoritmen."""
            digits = [int(d) for d in number]
            for i in range(len(digits) - 2, -1, -2):
                digits[i] *= 2
                if digits[i] > 9:
                    digits[i] -= 9
            return sum(digits) % 10 == 0
        
        if not luhn_check(clean):
            raise ValueError(
                f"Organisationsnummer '{v}' har ogiltig kontrollsiffra. "
                "Kontrollera att du skrivit rätt."
            )
        
        return clean


class OrgNummerInput(OrgNummerMixin):
    """Input för organisationsnummer med validering."""
    model_config = ConfigDict(extra="forbid")


class CompanyInfoInput(OrgNummerMixin):
    """Input för företagsinformation med validering."""
    model_config = ConfigDict(extra="forbid")
    response_format: ResponseFormat = Field(
        default=ResponseFormat.MARKDOWN,
        description="Svarsformat: markdown, json, csv, excel eller pdf"
    )


class FinansiellDataInput(OrgNummerMixin):
    """Input för finansiell data med validering."""
    model_config = ConfigDict(extra="forbid")
    index: int = Field(
        default=0,
        ge=0,
        le=10,
        description="Vilken årsredovisning (0=senaste, 1=näst senaste)"
    )
    response_format: ResponseFormat = Field(default=ResponseFormat.MARKDOWN)


class BatchInput(BaseModel):
    """Input för batch-sökning."""
    model_config = ConfigDict(extra="forbid")
    org_nummer_lista: List[str] = Field(
        min_length=1,
        max_length=20,
        description="Lista med organisationsnummer att söka (max 20)"
    )


class ExportInput(OrgNummerMixin):
    """Input för export med validering."""
    model_config = ConfigDict(extra="forbid")
    index: int = Field(default=0, ge=0)
    format: ResponseFormat = Field(
        default=ResponseFormat.PDF,
        description="Exportformat: pdf, excel, csv, json eller markdown"
    )
    filename: Optional[str] = Field(
        default=None,
        description="Valfritt filnamn (genereras automatiskt om ej angivet)"
    )


# =============================================================================
# FÖRBÄTTRING #2 & #3: Resources (passiv data via URI-schema)
# =============================================================================

@mcp.resource("bolagsverket://company/{org_nummer}")
def resource_company(org_nummer: str) -> str:
    """
    Företagsinformation som resurs.
    
    URI: bolagsverket://company/{org_nummer}
    Exempel: bolagsverket://company/5567671267
    """
    try:
        logger.info(f"Resource request: company/{org_nummer}")
        valid, clean_nr = validate_org_nummer(org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, clean_nr, org_nummer=org_nummer)
        
        info = fetch_company_info(clean_nr)
        return export_to_json(info)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e), org_nummer=org_nummer)


@mcp.resource("bolagsverket://financials/{org_nummer}")
def resource_financials(org_nummer: str) -> str:
    """
    Finansiell data (nyckeltal) som resurs.
    
    URI: bolagsverket://financials/{org_nummer}
    """
    try:
        logger.info(f"Resource request: financials/{org_nummer}")
        arsred, _, _ = fetch_and_parse_arsredovisning(org_nummer, 0)
        return export_to_json(arsred.nyckeltal)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e), org_nummer=org_nummer)


@mcp.resource("bolagsverket://people/{org_nummer}")
def resource_people(org_nummer: str) -> str:
    """
    Personer kopplade till företaget som resurs.
    
    URI: bolagsverket://people/{org_nummer}
    """
    try:
        logger.info(f"Resource request: people/{org_nummer}")
        arsred, _, _ = fetch_and_parse_arsredovisning(org_nummer, 0)
        return export_to_json([asdict(p) for p in arsred.personer])
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e), org_nummer=org_nummer)


@mcp.resource("bolagsverket://annual-reports/{org_nummer}")
def resource_annual_reports_list(org_nummer: str) -> str:
    """
    Lista tillgängliga årsredovisningar som resurs.
    
    URI: bolagsverket://annual-reports/{org_nummer}
    """
    try:
        logger.info(f"Resource request: annual-reports/{org_nummer}")
        clean_nr = clean_org_nummer(org_nummer)
        dok_data = make_api_request("POST", "/dokumentlista", {"identitetsbeteckning": clean_nr})
        dokument = dok_data.get("dokument", [])
        
        result = []
        for i, dok in enumerate(dokument):
            result.append({
                "index": i,
                "dokument_id": dok.get("dokumentId"),
                "period_fran": dok.get("rakenskapsperiod", {}).get("fran"),
                "period_till": dok.get("rakenskapsperiod", {}).get("till"),
                "inlamningsdatum": dok.get("inlamningsdatum"),
            })
        
        return export_to_json({"org_nummer": format_org_nummer(clean_nr), "arsredovisningar": result})
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e), org_nummer=org_nummer)


# =============================================================================
# NYA RESOURCES I v5.0
# =============================================================================

@mcp.resource("bolagsverket://risk/{org_nummer}")
def resource_risk_assessment(org_nummer: str) -> str:
    """
    Riskbedömning och röda flaggor som resurs.
    
    URI: bolagsverket://risk/{org_nummer}
    
    Returnerar:
    - Identifierade röda flaggor
    - Allvarlighetsgrad (kritisk/varning/info)
    - Rekommendationer
    """
    try:
        logger.info(f"Resource request: risk/{org_nummer}")
        valid, clean_nr = validate_org_nummer(org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, clean_nr, org_nummer=org_nummer)
        
        # Hämta senaste årsredovisningen och analysera
        arsred, _, _ = fetch_and_parse_arsredovisning(clean_nr, 0)
        full = parse_arsredovisning_full(clean_nr, 0)
        
        result = {
            "org_nummer": format_org_nummer(clean_nr),
            "foretag_namn": full.foretag_namn,
            "antal_flaggor": len(full.roda_flaggor),
            "kritiska": [asdict(f) for f in full.roda_flaggor if f.allvarlighet == "kritisk"],
            "varningar": [asdict(f) for f in full.roda_flaggor if f.allvarlighet == "varning"],
            "info": [asdict(f) for f in full.roda_flaggor if f.allvarlighet == "info"],
        }
        
        return export_to_json(result)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e), org_nummer=org_nummer)


@mcp.resource("bolagsverket://nyckeltal/{org_nummer}/{year}")
def resource_nyckeltal_year(org_nummer: str, year: str) -> str:
    """
    Nyckeltal för specifikt räkenskapsår som resurs.
    
    URI: bolagsverket://nyckeltal/{org_nummer}/{year}
    Exempel: bolagsverket://nyckeltal/5567671267/2023
    
    Parametrar:
    - org_nummer: Organisationsnummer
    - year: Räkenskapsår (t.ex. "2023")
    """
    try:
        logger.info(f"Resource request: nyckeltal/{org_nummer}/{year}")
        valid, clean_nr = validate_org_nummer(org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, clean_nr, org_nummer=org_nummer)
        
        # Hitta rätt årsredovisning baserat på år
        dok_data = make_api_request("POST", "/dokumentlista", {"identitetsbeteckning": clean_nr})
        dokument = dok_data.get("dokument", [])
        
        target_year = int(year)
        found_index = None
        
        for i, dok in enumerate(dokument):
            period_till = dok.get("rakenskapsperiod", {}).get("till", "")
            if period_till and period_till.startswith(str(target_year)):
                found_index = i
                break
        
        if found_index is None:
            return handle_error(
                ErrorCode.ANNUAL_REPORT_NOT_FOUND, 
                f"Ingen årsredovisning hittades för år {year}",
                org_nummer=org_nummer,
                year=year
            )
        
        arsred, _, _ = fetch_and_parse_arsredovisning(clean_nr, found_index)
        nyckeltal = arsred.nyckeltal
        
        result = {
            "org_nummer": format_org_nummer(clean_nr),
            "rakenskapsar": year,
            "nettoomsattning": nyckeltal.nettoomsattning,
            "arets_resultat": nyckeltal.arets_resultat,
            "eget_kapital": nyckeltal.eget_kapital,
            "balansomslutning": nyckeltal.balansomslutning,
            "soliditet": nyckeltal.soliditet,
            "vinstmarginal": nyckeltal.vinstmarginal,
            "roe": nyckeltal.roe,
            "antal_anstallda": nyckeltal.antal_anstallda,
        }
        
        return export_to_json(result)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e), org_nummer=org_nummer)


@mcp.resource("bolagsverket://status/{org_nummer}")
def resource_company_status(org_nummer: str) -> str:
    """
    Företagets aktuella status som resurs.
    
    URI: bolagsverket://status/{org_nummer}
    
    Returnerar:
    - Aktiv/Avregistrerad
    - Pågående konkurs/likvidation
    - Avregistreringsorsak (om tillämpligt)
    """
    try:
        logger.info(f"Resource request: status/{org_nummer}")
        valid, clean_nr = validate_org_nummer(org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, clean_nr, org_nummer=org_nummer)
        
        info = fetch_company_info(clean_nr)
        
        result = {
            "org_nummer": format_org_nummer(clean_nr),
            "namn": info.namn,
            "status": info.status,
            "verksam": info.verksam_organisation,
            "pagaende_konkurs": info.pagaende_konkurs,
            "pagaende_likvidation": info.pagaende_likvidation,
            "avregistreringsdatum": info.avregistreringsdatum,
            "avregistreringsorsak": info.avregistreringsorsak,
            "reklamsparr": info.reklamsparr,
        }
        
        return export_to_json(result)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e), org_nummer=org_nummer)


@mcp.resource("bolagsverket://bransch/{org_nummer}")
def resource_bransch(org_nummer: str) -> str:
    """
    Företagets branschklassificering (SNI-koder) som resurs.
    
    URI: bolagsverket://bransch/{org_nummer}
    
    Returnerar:
    - SNI-koder med klartext
    - Verksamhetsbeskrivning
    """
    try:
        logger.info(f"Resource request: bransch/{org_nummer}")
        valid, clean_nr = validate_org_nummer(org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, clean_nr, org_nummer=org_nummer)
        
        info = fetch_company_info(clean_nr)
        
        result = {
            "org_nummer": format_org_nummer(clean_nr),
            "namn": info.namn,
            "sni_koder": info.sni_koder,
            "verksamhetsbeskrivning": info.verksamhet,
        }
        
        return export_to_json(result)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e), org_nummer=org_nummer)


@mcp.resource("bolagsverket://server-info")
def resource_server_info() -> str:
    """
    Information om MCP-servern.
    
    URI: bolagsverket://server-info
    
    Returnerar:
    - Version
    - Tillgängliga verktyg
    - Tillgängliga resurser
    - API-information
    """
    return export_to_json({
        "version": "5.0.1",
        "namn": "Bolagsverket MCP Server",
        "api": {
            "namn": "Värdefulla datamängder",
            "url": "https://bolagsverket.se/apierochoppnadata",
            "beskrivning": "API för företagsinformation och årsredovisningar"
        },
        "verktyg": [
            {"namn": "bolagsverket_analyze_full", "beskrivning": "Komplett årsredovisningsanalys"},
            {"namn": "bolagsverket_get_basic_info", "beskrivning": "Grundläggande företagsinfo"},
            {"namn": "bolagsverket_get_nyckeltal", "beskrivning": "Finansiella nyckeltal"},
            {"namn": "bolagsverket_get_styrelse", "beskrivning": "Styrelse och ledning"},
            {"namn": "bolagsverket_risk_check", "beskrivning": "Riskanalys och röda flaggor"},
            {"namn": "bolagsverket_trend", "beskrivning": "Historisk utveckling"},
            {"namn": "bolagsverket_search", "beskrivning": "Sök företag på namn"},
            {"namn": "bolagsverket_export_pro", "beskrivning": "Export till Word/PowerPoint"},
        ],
        "resurser": [
            {"uri": "bolagsverket://company/{org}", "beskrivning": "Företagsinfo"},
            {"uri": "bolagsverket://financials/{org}", "beskrivning": "Nyckeltal"},
            {"uri": "bolagsverket://nyckeltal/{org}/{år}", "beskrivning": "Nyckeltal för specifikt år"},
            {"uri": "bolagsverket://people/{org}", "beskrivning": "Personer"},
            {"uri": "bolagsverket://risk/{org}", "beskrivning": "Riskbedömning"},
            {"uri": "bolagsverket://status/{org}", "beskrivning": "Status"},
            {"uri": "bolagsverket://bransch/{org}", "beskrivning": "Branschinfo"},
            {"uri": "bolagsverket://annual-reports/{org}", "beskrivning": "Årsredovisningar"},
        ],
        "taxonomier_stöd": {
            "aktuella": [
                "K2 Aktiebolag 2024-09-12",
                "K2 Aktiebolag 2021-10-31",
                "K2 Ek.för/HB/KB/Filial 2024-09-12",
                "K3 Aktiebolag 2021-10-31",
                "K3K Koncern 2021-10-31",
                "Revisionsberättelse 2020-12-01",
                "Fastställelseintyg 2020-12-01"
            ],
            "arkiverade": [
                "K2 2017-09-30 (stöds EJ av Bolagsverket sedan 2023)",
                "K3 2018-12-17 (stöds EJ av Bolagsverket sedan 2023)"
            ]
        }
    })


# =============================================================================
# FÖRBÄTTRING #4: Prompts (fördefinierade arbetsflöden)
# =============================================================================

@mcp.prompt("due-diligence")
def prompt_due_diligence(org_nummer: str) -> str:
    """
    Due diligence-analys av ett företag.
    
    Utför en komplett genomlysning med:
    - Grundläggande företagsinformation
    - Finansiell analys med nyckeltal
    - Identifiering av nyckelpersoner
    - Riskbedömning
    """
    return f"""Utför en due diligence-analys av företaget med organisationsnummer {org_nummer}.

## Steg 1: Grundläggande information
Använd verktyget `bolagsverket_get_basic_info` för att hämta:
- Företagsnamn och organisationsform
- Registreringsdatum och status
- Säte och adress
- Verksamhetsbeskrivning
- SNI-koder (bransch)

## Steg 2: Finansiell analys
Använd verktyget `bolagsverket_get_nyckeltal` för att analysera:
- Omsättning och resultat
- Soliditet och likviditet
- Vinstmarginal och ROE
- Jämför med föregående år om möjligt

## Steg 3: Nyckelpersoner
Använd verktyget `bolagsverket_get_styrelse` för att identifiera:
- VD och företrädare
- Styrelseledamöter
- Revisorer

## Steg 4: Trendanalys
Använd verktyget `bolagsverket_get_trends` för att:
- Analysera omsättningstillväxt
- Identifiera resultatutveckling
- Bedöma finansiell stabilitet över tid

## Steg 5: Sammanfattning
Ge en övergripande bedömning med:
- Styrkor
- Svagheter/risker
- Rekommendation
"""


@mcp.prompt("compare-companies")
def prompt_compare_companies(org_nummer_1: str, org_nummer_2: str) -> str:
    """
    Jämför två företag.
    
    Skapar en side-by-side jämförelse av nyckeltal och egenskaper.
    """
    return f"""Jämför följande två företag:

**Företag 1:** {org_nummer_1}
**Företag 2:** {org_nummer_2}

## Steg 1: Hämta grundinformation
Använd `bolagsverket_get_basic_info` för båda företagen.

## Steg 2: Hämta nyckeltal
Använd `bolagsverket_get_nyckeltal` för båda företagen.

## Steg 3: Skapa jämförelsetabell
Presentera en tabell med följande jämförelsepunkter:
- Omsättning
- Årets resultat
- Eget kapital
- Soliditet
- Vinstmarginal
- Antal anställda
- Organisationsform
- Bransch (SNI)

## Steg 4: Analys
- Vilket företag är större?
- Vilket har bättre lönsamhet?
- Vilket har starkare finansiell ställning?
- Eventuella risker eller fördelar
"""


@mcp.prompt("person-network")
def prompt_person_network(org_nummer: str) -> str:
    """
    Analysera personkopplingar för ett företag.
    
    Identifierar nyckelpersoner och deras roller.
    """
    return f"""Analysera nyckelpersoner kopplade till företaget {org_nummer}.

## Steg 1: Identifiera personer
Använd `bolagsverket_get_styrelse` för att hämta:
- VD
- Styrelseledamöter
- Revisorer

## Steg 2: Analysera roller
För varje person, notera:
- Fullständigt namn
- Roll i företaget
- Om de har flera roller

## Steg 3: Sammanfattning
Presentera en översikt av företagets ledning och governance-struktur.
"""


@mcp.prompt("export-report")
def prompt_export_report(org_nummer: str, format: str = "pdf") -> str:
    """
    Exportera en komplett företagsrapport.
    """
    return f"""Skapa och exportera en komplett rapport för företaget {org_nummer}.

## Steg 1: Samla data
Hämta all relevant information med:
- `bolagsverket_get_basic_info`
- `bolagsverket_get_nyckeltal`
- `bolagsverket_get_styrelse`

## Steg 2: Exportera
Använd `bolagsverket_export` med format="{format}" för att skapa en nedladdningsbar fil.

## Steg 3: Presentera
Ge användaren länk till den exporterade filen.
"""


# =============================================================================
# NYA PROMPTS I v5.0
# =============================================================================

@mcp.prompt("quick-check")
def prompt_quick_check(org_nummer: str) -> str:
    """
    Snabb kontroll av ett företag.
    
    Perfekt för att snabbt verifiera ett företags status och grundläggande info.
    """
    return f"""Gör en snabb kontroll av företaget {org_nummer}.

## Snabbkontroll
1. Använd `bolagsverket_get_basic_info` för att verifiera:
   - Att företaget existerar
   - Att det är aktivt (inte avregistrerat)
   - Organisationsform

2. Använd `bolagsverket_risk_check` för att snabbt identifiera eventuella:
   - Röda flaggor
   - Pågående konkurs/likvidation
   - Kritiska varningar

## Sammanfattning
Ge ett kort svar (max 3 meningar) om företagets status:
- ✅ OK - om allt ser bra ut
- ⚠️ Varning - om det finns saker att vara uppmärksam på
- 🚨 Risk - om det finns allvarliga problem
"""


@mcp.prompt("investment-analysis")
def prompt_investment_analysis(org_nummer: str) -> str:
    """
    Investeringsanalys av ett företag.
    
    Djupgående analys för potentiella investerare.
    """
    return f"""Utför en investeringsanalys av företaget {org_nummer}.

## Del 1: Företagsöversikt
Använd `bolagsverket_get_basic_info` och `bolagsverket_get_verksamhet` för att förstå:
- Vad gör företaget?
- Hur länge har det funnits?
- Vilken bransch verkar det i?

## Del 2: Finansiell hälsa
Använd `bolagsverket_get_nyckeltal` för att analysera:
- **Tillväxt**: Ökar eller minskar omsättningen?
- **Lönsamhet**: Vinstmarginal och ROE
- **Stabilitet**: Soliditet och eget kapital
- **Storlek**: Omsättning och antal anställda

## Del 3: Trendanalys
Använd `bolagsverket_trend` för att se:
- Hur har omsättningen utvecklats de senaste åren?
- Är resultatet stabilt eller volatilt?
- Finns det tecken på tillväxt eller nedgång?

## Del 4: Riskbedömning
Använd `bolagsverket_risk_check` för att identifiera:
- Finansiella risker
- Röda flaggor
- Varningar

## Del 5: Ledning
Använd `bolagsverket_get_styrelse` för att se:
- Vem leder företaget?
- Erfarenhet i styrelsen?

## Investeringsrekommendation
Sammanfatta med:
- **Styrkor** (3 punkter)
- **Risker** (3 punkter)
- **Rekommendation**: Köp / Avvakta / Undvik
"""


@mcp.prompt("supplier-check")
def prompt_supplier_check(org_nummer: str) -> str:
    """
    Leverantörskontroll.
    
    Perfekt för att verifiera en potentiell leverantör innan avtal.
    """
    return f"""Utför en leverantörskontroll av företaget {org_nummer}.

## Kontrollpunkter

### 1. Grundläggande verifiering
Använd `bolagsverket_get_basic_info`:
- [ ] Företaget existerar och är aktivt
- [ ] Organisationsform är lämplig
- [ ] Har funnits i minst 2 år

### 2. Finansiell stabilitet
Använd `bolagsverket_get_nyckeltal`:
- [ ] Positivt eget kapital
- [ ] Soliditet över 20%
- [ ] Positivt resultat

### 3. Riskfaktorer
Använd `bolagsverket_risk_check`:
- [ ] Ingen pågående konkurs
- [ ] Ingen pågående likvidation
- [ ] Inga kritiska röda flaggor

### 4. Historik
Använd `bolagsverket_trend`:
- [ ] Stabil eller växande omsättning
- [ ] Inga dramatiska nedgångar

## Sammanfattning
Ange en av följande bedömningar:
- ✅ **GODKÄND** - Företaget verkar stabilt och pålitligt
- ⚠️ **GODKÄND MED RESERVATION** - Finns några frågetecken, men acceptabelt
- 🚨 **EJ REKOMMENDERAD** - Betydande risker identifierade
"""


@mcp.prompt("batch-analysis")
def prompt_batch_analysis(org_nummer_lista: str) -> str:
    """
    Analysera flera företag samtidigt.
    
    Ange organisationsnummer separerade med komma.
    """
    org_list = [o.strip() for o in org_nummer_lista.split(",")]
    org_formatted = "\n".join([f"- {org}" for org in org_list])
    
    return f"""Analysera följande företag och skapa en jämförelsetabell:

{org_formatted}

## Steg 1: Samla data
För varje företag, hämta:
- Grundinfo (`bolagsverket_get_basic_info`)
- Nyckeltal (`bolagsverket_get_nyckeltal`)

## Steg 2: Skapa jämförelsetabell
| Företag | Omsättning | Resultat | Soliditet | Anställda |
|---------|------------|----------|-----------|-----------|
| ... | ... | ... | ... | ... |

## Steg 3: Rangordna
Rangordna företagen efter:
1. Finansiell styrka
2. Storlek
3. Lönsamhet

## Steg 4: Sammanfattning
Ge en kort analys av:
- Vilket företag är starkast?
- Finns det några som sticker ut negativt?
"""


@mcp.prompt("annual-report-summary")
def prompt_annual_report_summary(org_nummer: str) -> str:
    """
    Sammanfatta senaste årsredovisningen.
    
    Ger en lättläst sammanfattning av årsredovisningen.
    """
    return f"""Skapa en lättläst sammanfattning av senaste årsredovisningen för {org_nummer}.

## Använd verktyg
Använd `bolagsverket_analyze_full` för att hämta all data.

## Struktur för sammanfattningen

### 🏢 Om företaget
- Namn och organisationsform
- Vad gör de? (verksamhetsbeskrivning)

### 📊 Finansiellt (i klarspråk)
- "Företaget omsatte X miljoner kronor" 
- "Vinsten/förlusten blev Y kronor"
- "Företaget har Z kronor i eget kapital"

### 👥 Ledning
- VD: [namn]
- Styrelseordförande: [namn]

### 📈 Hur går det?
Jämför med föregående år:
- Bättre eller sämre omsättning?
- Bättre eller sämre resultat?

### ⚠️ Att vara medveten om
Lista eventuella röda flaggor eller varningar.

### 💡 Slutsats
En mening som sammanfattar företagets situation.
"""


# =============================================================================
# FÖRBÄTTRING #9: Granulära verktyg (Single Responsibility)
# =============================================================================

@mcp.tool()
def bolagsverket_check_status() -> str:
    """
    Kontrollera om Bolagsverkets API är tillgängligt.
    
    Returnerar API:ets status.
    """
    try:
        logger.info("Kontrollerar API-status...")
        token = token_manager.get_token()
        with httpx.Client(timeout=30.0) as client:
            response = client.get(f"{BASE_URL}/isalive", headers={"Authorization": f"Bearer {token}"})
        if response.status_code == 200:
            return "✅ Bolagsverkets API är tillgängligt!"
        return f"⚠️ API svarade med status {response.status_code}"
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


@mcp.tool()
def bolagsverket_get_basic_info(params: OrgNummerInput) -> str:
    """
    Hämta grundläggande företagsinformation.
    
    Returnerar:
    - Företagsnamn
    - Organisationsform (AB, HB, etc.)
    - Juridisk form
    - Registreringsdatum
    - Status (aktiv/avregistrerad)
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result, org_nummer=params.org_nummer)
        
        info = fetch_company_info(result)
        
        lines = [
            f"# {info.namn}",
            f"",
            f"**Organisationsnummer:** {info.org_nummer}",
            f"**Organisationsform:** {info.organisationsform}",
        ]
        
        if info.juridisk_form:
            lines.append(f"**Juridisk form:** {info.juridisk_form}")
        
        lines.append(f"**Registreringsdatum:** {info.registreringsdatum}")
        lines.append(f"**Status:** {info.status}")
        
        if info.avregistreringsdatum:
            lines.append(f"**Avregistreringsdatum:** {info.avregistreringsdatum[:10]}")
        
        return "\n".join(lines)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e), org_nummer=params.org_nummer)


@mcp.tool()
def bolagsverket_get_address(params: OrgNummerInput) -> str:
    """
    Hämta företagets postadress.
    
    Returnerar:
    - Utdelningsadress
    - Postnummer
    - Postort
    - Säte (län)
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        info = fetch_company_info(result)
        
        lines = [
            f"# Adress: {info.namn}",
            f"",
        ]
        
        if info.adress.get('utdelningsadress'):
            lines.append(info.adress['utdelningsadress'])
        
        lines.append(f"{info.adress.get('postnummer', '')} {info.adress.get('postort', '')}")
        
        if info.sate:
            lines.append(f"**Säte:** {info.sate}")
        
        return "\n".join(lines)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


@mcp.tool()
def bolagsverket_get_verksamhet(params: OrgNummerInput) -> str:
    """
    Hämta företagets verksamhetsbeskrivning och branschkoder.
    
    Returnerar:
    - Verksamhetsbeskrivning
    - SNI-koder med klartext
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        info = fetch_company_info(result)
        
        lines = [
            f"# Verksamhet: {info.namn}",
            f"",
        ]
        
        if info.verksamhet:
            lines.append("## Beskrivning")
            lines.append(info.verksamhet.strip())
            lines.append("")
        
        if info.sni_koder:
            lines.append("## SNI-koder (bransch)")
            for sni in info.sni_koder:
                lines.append(f"- **{sni['kod']}**: {sni['klartext']}")
        else:
            lines.append("*Inga SNI-koder registrerade*")
        
        return "\n".join(lines)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


@mcp.tool()
def bolagsverket_get_nyckeltal(params: FinansiellDataInput) -> str:
    """
    Hämta finansiella nyckeltal från årsredovisning.
    
    Returnerar:
    - Nettoomsättning
    - Resultat efter finansiella poster
    - Årets resultat
    - Eget kapital
    - Soliditet
    - Vinstmarginal (beräknad)
    - ROE (beräknad)
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        arsredovisning, _, _ = fetch_and_parse_arsredovisning(result, params.index)
        
        if params.response_format == ResponseFormat.JSON:
            return export_to_json(arsredovisning.nyckeltal)
        
        nyckeltal = arsredovisning.nyckeltal
        lines = [
            f"# Nyckeltal: {arsredovisning.foretag_namn}",
            f"**Räkenskapsår:** {arsredovisning.rakenskapsar_start} – {arsredovisning.rakenskapsar_slut}",
            "",
            "| Nyckeltal | Värde |",
            "|-----------|------:|"
        ]
        
        data = [
            ("Nettoomsättning", nyckeltal.nettoomsattning, "SEK"),
            ("Resultat efter fin. poster", nyckeltal.resultat_efter_finansiella, "SEK"),
            ("Årets resultat", nyckeltal.arets_resultat, "SEK"),
            ("Eget kapital", nyckeltal.eget_kapital, "SEK"),
            ("Balansomslutning", nyckeltal.balansomslutning, "SEK"),
            ("Soliditet", nyckeltal.soliditet, "%"),
            ("Vinstmarginal", nyckeltal.vinstmarginal, "%"),
            ("ROE", nyckeltal.roe, "%"),
            ("Antal anställda", nyckeltal.antal_anstallda, "st"),
        ]
        
        for label, value, unit in data:
            if value is not None:
                if isinstance(value, int) and unit == "SEK":
                    lines.append(f"| {label} | {value:,} {unit} |")
                else:
                    lines.append(f"| {label} | {value} {unit} |")
        
        return "\n".join(lines)
    except Exception as e:
        return handle_error(ErrorCode.ANNUAL_REPORT_NOT_FOUND, str(e), org_nummer=params.org_nummer)


@mcp.tool()
def bolagsverket_get_styrelse(params: OrgNummerInput) -> str:
    """
    Hämta styrelse, VD och revisorer från årsredovisning.
    
    Returnerar:
    - VD (Verkställande direktör)
    - Styrelseledamöter
    - Revisorer
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        arsredovisning, _, _ = fetch_and_parse_arsredovisning(result, 0)
        
        lines = [
            f"# Personer: {arsredovisning.foretag_namn}",
            f"**Org.nr:** {format_org_nummer(arsredovisning.org_nummer)}",
            f"**Källa:** Årsredovisning {arsredovisning.rakenskapsar_slut[:4]}",
            "",
            "| Namn | Roll |",
            "|------|------|",
        ]
        
        for person in arsredovisning.personer:
            lines.append(f"| {person.fullnamn} | {person.roll} |")
        
        if not arsredovisning.personer:
            lines.append("| *Inga personer hittades* | - |")
        
        return "\n".join(lines)
    except Exception as e:
        return handle_error(ErrorCode.ANNUAL_REPORT_NOT_FOUND, str(e), org_nummer=params.org_nummer)


@mcp.tool()
def bolagsverket_get_trends(params: OrgNummerInput) -> str:
    """
    Flerårsöversikt med trendanalys (upp till 4 år).
    
    Returnerar:
    - Nyckeltal per år
    - Omsättningstillväxt (%)
    - Resultatutveckling
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        clean_nr = clean_org_nummer(result)
        
        dok_data = make_api_request("POST", "/dokumentlista", {"identitetsbeteckning": clean_nr})
        dokument = dok_data.get("dokument", [])
        
        if not dokument:
            return handle_error(ErrorCode.ANNUAL_REPORT_NOT_FOUND, "Inga årsredovisningar hittades")
        
        dok_id = dokument[0].get("dokumentId")
        zip_bytes = download_document_bytes(dok_id)
        
        xhtml_content = None
        with zipfile.ZipFile(BytesIO(zip_bytes)) as zf:
            for name in zf.namelist():
                if name.lower().endswith(('.xhtml', '.html')):
                    xhtml_content = zf.read(name).decode('utf-8')
                    break
        
        if not xhtml_content:
            return handle_error(ErrorCode.PARSE_ERROR, "Ingen XHTML hittades")
        
        parser = IXBRLParser(xhtml_content)
        metadata = parser.get_metadata()
        oversikt = parser.get_flerarsoversikt()
        
        lines = [
            f"# Flerårsöversikt: {metadata['foretag_namn']}",
            f"**Org.nr:** {format_org_nummer(metadata['org_nummer'])}",
            "",
        ]
        
        periods = sorted(oversikt.keys())
        if not periods:
            return handle_error(ErrorCode.PARSE_ERROR, "Kunde inte extrahera flerårsdata")
        
        header = "| Nyckeltal |"
        separator = "|-----------|"
        for p in periods:
            year = f"År {p[-1]}"
            header += f" {year} |"
            separator += "------:|"
        lines.append(header)
        lines.append(separator)
        
        metrics = [
            ('nettoomsattning', 'Nettoomsättning', True),
            ('resultat_efter_finansiella', 'Resultat', True),
            ('arets_resultat', 'Årets resultat', True),
            ('eget_kapital', 'Eget kapital', True),
            ('soliditet', 'Soliditet (%)', False),
        ]
        
        for key, label, is_money in metrics:
            row = f"| {label} |"
            for p in periods:
                val = getattr(oversikt[p], key, None)
                if val is not None:
                    if is_money:
                        row += f" {val:,} |"
                    else:
                        row += f" {val} |"
                else:
                    row += " - |"
            lines.append(row)
        
        if len(periods) >= 2:
            lines.append("")
            lines.append("## Tillväxt")
            
            first = oversikt[periods[-1]]
            last = oversikt[periods[0]]
            
            if first.nettoomsattning and last.nettoomsattning and first.nettoomsattning > 0:
                tillvaxt = ((last.nettoomsattning - first.nettoomsattning) / first.nettoomsattning) * 100
                lines.append(f"- **Omsättningstillväxt:** {tillvaxt:.1f}%")
            
            if first.arets_resultat and last.arets_resultat:
                diff = last.arets_resultat - first.arets_resultat
                lines.append(f"- **Resultatförändring:** {diff:+,} SEK")
        
        return "\n".join(lines)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


@mcp.tool()
def bolagsverket_batch_lookup(params: BatchInput) -> str:
    """
    Sök information om flera företag samtidigt (max 20).
    
    Returnerar sammanställd tabell med:
    - Företagsnamn
    - Status
    - Omsättning
    - Resultat
    """
    try:
        results = []
        errors = []
        
        for org_nr in params.org_nummer_lista:
            try:
                clean_nr = clean_org_nummer(org_nr)
                
                org_data = make_api_request("POST", "/organisationer", {"identitetsbeteckning": clean_nr})
                orgs = org_data.get("organisationer", [])
                
                if not orgs:
                    errors.append(f"{org_nr}: Hittades ej")
                    continue
                
                org = orgs[0]
                namn = org.get("organisationsnamn", {}).get("organisationsnamnLista", [{}])[0].get("namn", "Okänt")
                form = org.get("organisationsform", {}).get("klartext", "-")
                
                avreg = org.get("avregistreradOrganisation", {})
                status = "Avregistrerad" if avreg and avreg.get("avregistreringsdatum") else "Aktiv"
                
                nyckeltal = None
                try:
                    arsred, _, _ = fetch_and_parse_arsredovisning(org_nr, 0)
                    nyckeltal = arsred.nyckeltal
                except:
                    pass
                
                results.append({
                    'org_nr': format_org_nummer(clean_nr),
                    'namn': namn,
                    'form': form,
                    'status': status,
                    'nyckeltal': nyckeltal
                })
                
            except Exception as e:
                errors.append(f"{org_nr}: {str(e)}")
        
        lines = [
            f"# Batch-sökning ({len(results)} företag)",
            "",
            "| Org.nr | Företag | Form | Status | Omsättning | Resultat |",
            "|--------|---------|------|--------|------------|----------|",
        ]
        
        for r in results:
            oms = f"{r['nyckeltal'].nettoomsattning:,}" if r['nyckeltal'] and r['nyckeltal'].nettoomsattning else "-"
            res = f"{r['nyckeltal'].arets_resultat:,}" if r['nyckeltal'] and r['nyckeltal'].arets_resultat else "-"
            lines.append(f"| {r['org_nr']} | {r['namn'][:25]} | {r['form'][:10]} | {r['status']} | {oms} | {res} |")
        
        if errors:
            lines.append("")
            lines.append("## Fel")
            for e in errors:
                lines.append(f"- {e}")
        
        return "\n".join(lines)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


@mcp.tool()
def bolagsverket_export(params: ExportInput) -> str:
    """
    Exportera årsredovisningsdata till fil.
    
    Tillgängliga format:
    - pdf: Snygg PDF-rapport
    - excel: Excel med formatering
    - csv: CSV för import
    - json: Strukturerad JSON
    - markdown: Markdown-text
    - xhtml: Original iXBRL-fil från Bolagsverket
    - zip: Original ZIP-arkiv från Bolagsverket
    
    Filer sparas i ~/Downloads/bolagsverket/
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        arsredovisning, xhtml_bytes, zip_bytes = fetch_and_parse_arsredovisning(result, params.index)
        
        # Generera filnamn
        clean_name = re.sub(r'[^\w\s-]', '', arsredovisning.foretag_namn).strip()
        year = arsredovisning.rakenskapsar_slut[:4] if arsredovisning.rakenskapsar_slut else "unknown"
        
        if params.format == ResponseFormat.ZIP:
            filename = params.filename or f"{clean_name}_{year}_arsredovisning.zip"
            filepath = ensure_output_dir() / filename
            with open(filepath, 'wb') as f:
                f.write(zip_bytes)
            logger.info(f"ZIP exporterad till: {filepath}")
            return f"✅ Original ZIP exporterad till: {filepath}"
        
        elif params.format == ResponseFormat.XHTML:
            filename = params.filename or f"{clean_name}_{year}_arsredovisning.xhtml"
            filepath = ensure_output_dir() / filename
            with open(filepath, 'wb') as f:
                f.write(xhtml_bytes)
            logger.info(f"XHTML exporterad till: {filepath}")
            return f"✅ Original XHTML (iXBRL) exporterad till: {filepath}"
        
        elif params.format == ResponseFormat.PDF:
            filepath = export_to_pdf(arsredovisning, params.filename)
            return f"✅ PDF exporterad till: {filepath}"
        
        elif params.format == ResponseFormat.EXCEL:
            filepath = export_to_excel(arsredovisning, params.filename)
            return f"✅ Excel exporterad till: {filepath}"
        
        elif params.format == ResponseFormat.CSV:
            filename = params.filename or f"{arsredovisning.foretag_namn}.csv"
            filepath = export_to_csv(asdict(arsredovisning.nyckeltal), filename)
            return f"✅ CSV exporterad till: {filepath}"
        
        elif params.format == ResponseFormat.JSON:
            return export_to_json(arsredovisning)
        
        else:
            return export_to_markdown(arsredovisning)
        
    except Exception as e:
        return handle_error(ErrorCode.EXPORT_ERROR, str(e), format=params.format.value)


@mcp.tool()
def bolagsverket_list_arsredovisningar(params: OrgNummerInput) -> str:
    """
    Lista tillgängliga årsredovisningar för ett företag.
    
    Returnerar:
    - Räkenskapsperiod
    - Inlämningsdatum
    - Dokument-ID
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        clean_nr = clean_org_nummer(result)
        dok_data = make_api_request("POST", "/dokumentlista", {"identitetsbeteckning": clean_nr})
        dokument = dok_data.get("dokument", [])
        
        if not dokument:
            return f"Inga årsredovisningar hittades för {format_org_nummer(clean_nr)}"
        
        lines = [
            f"# Årsredovisningar: {format_org_nummer(clean_nr)}",
            "",
            "| # | Period | Inlämnad |",
            "|---|--------|----------|",
        ]
        
        for i, dok in enumerate(dokument):
            period = dok.get("rakenskapsperiod", {})
            fran = period.get("fran", "-")
            till = period.get("till", "-")
            inlamnad = dok.get("inlamningsdatum", "-")
            lines.append(f"| {i} | {fran} – {till} | {inlamnad} |")
        
        lines.append("")
        lines.append(f"*Använd index 0-{len(dokument)-1} för att hämta specifik årsredovisning*")
        
        return "\n".join(lines)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


class DownloadInput(BaseModel):
    """Input för nedladdning av original-årsredovisning."""
    model_config = ConfigDict(extra="forbid")
    org_nummer: str = Field(min_length=10, description="Organisationsnummer")
    index: int = Field(default=0, ge=0, description="Vilken årsredovisning (0=senaste)")
    format: str = Field(
        default="zip",
        description="Format: 'zip' för original ZIP-arkiv, 'xhtml' för iXBRL-fil"
    )
    destination: Optional[str] = Field(
        default=None,
        description="Valfri destination (t.ex. ~/Desktop). Standard: ~/Downloads/bolagsverket/"
    )


@mcp.tool()
def bolagsverket_download_original(params: DownloadInput) -> str:
    """
    Ladda ner original-årsredovisning från Bolagsverket.
    
    Detta verktyg sparar den RIKTIGA årsredovisningen som Bolagsverket tillhandahåller,
    inte en genererad sammanfattning.
    
    Format:
    - zip: Original ZIP-arkiv som innehåller iXBRL-filen
    - xhtml: Extraherad iXBRL/XHTML-fil (kan öppnas i webbläsare)
    
    Returnerar sökväg till den sparade filen.
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        arsredovisning, xhtml_bytes, zip_bytes = fetch_and_parse_arsredovisning(result, params.index)
        
        # Generera filnamn
        clean_name = re.sub(r'[^\w\s-]', '', arsredovisning.foretag_namn).strip().replace(' ', '_')
        year = arsredovisning.rakenskapsar_slut[:4] if arsredovisning.rakenskapsar_slut else "unknown"
        
        # Bestäm destination
        if params.destination:
            dest_path = Path(params.destination).expanduser()
            dest_path.mkdir(parents=True, exist_ok=True)
        else:
            dest_path = ensure_output_dir()
        
        if params.format.lower() == "zip":
            filename = f"{clean_name}_{year}_arsredovisning.zip"
            filepath = dest_path / filename
            with open(filepath, 'wb') as f:
                f.write(zip_bytes)
            size_kb = len(zip_bytes) / 1024
            logger.info(f"Original ZIP sparad: {filepath} ({size_kb:.1f} KB)")
            return f"✅ Original årsredovisning (ZIP) sparad:\n📁 {filepath}\n📊 Storlek: {size_kb:.1f} KB\n📅 Räkenskapsår: {arsredovisning.rakenskapsar_start} – {arsredovisning.rakenskapsar_slut}"
        
        elif params.format.lower() == "xhtml":
            filename = f"{clean_name}_{year}_arsredovisning.xhtml"
            filepath = dest_path / filename
            with open(filepath, 'wb') as f:
                f.write(xhtml_bytes)
            size_kb = len(xhtml_bytes) / 1024
            logger.info(f"Original XHTML sparad: {filepath} ({size_kb:.1f} KB)")
            return f"✅ Original årsredovisning (iXBRL/XHTML) sparad:\n📁 {filepath}\n📊 Storlek: {size_kb:.1f} KB\n📅 Räkenskapsår: {arsredovisning.rakenskapsar_start} – {arsredovisning.rakenskapsar_slut}\n💡 Tips: Öppna filen i en webbläsare för att se årsredovisningen"
        
        else:
            return handle_error(ErrorCode.INVALID_INPUT, f"Okänt format: {params.format}. Använd 'zip' eller 'xhtml'.")
        
    except Exception as e:
        return handle_error(ErrorCode.EXPORT_ERROR, str(e))


# =============================================================================
# NYA VERKTYG V4: Förbättring 1, 3, 4, 6, 7, 8, 9
# =============================================================================

# --- Hjälpfunktioner för nya verktyg ---

def analyze_roda_flaggor(parser, nyckeltal: Nyckeltal, nyckeltal_prev: Nyckeltal = None) -> List[RodFlagga]:
    """Analysera och identifiera röda flaggor."""
    flaggor = []
    
    # 1. Negativt eget kapital
    if nyckeltal.eget_kapital and nyckeltal.eget_kapital < 0:
        flaggor.append(RodFlagga(
            typ="negativt_eget_kapital",
            allvarlighet="kritisk",
            beskrivning=f"Negativt eget kapital: {nyckeltal.eget_kapital:,} SEK",
            varde=nyckeltal.eget_kapital,
            rekommendation="Bolaget kan behöva upprätta kontrollbalansräkning enligt ABL 25 kap."
        ))
    
    # 2. Kraftigt fallande omsättning
    if nyckeltal_prev and nyckeltal.nettoomsattning and nyckeltal_prev.nettoomsattning:
        if nyckeltal_prev.nettoomsattning > 0:
            forandring = ((nyckeltal.nettoomsattning - nyckeltal_prev.nettoomsattning) / nyckeltal_prev.nettoomsattning) * 100
            if forandring < -20:
                flaggor.append(RodFlagga(
                    typ="fallande_omsattning",
                    allvarlighet="varning",
                    beskrivning=f"Omsättningen har minskat med {abs(forandring):.1f}%",
                    varde=forandring,
                    rekommendation="Analysera orsakerna till omsättningsminskningen"
                ))
    
    # 3. Låg/negativ soliditet
    if nyckeltal.soliditet is not None:
        if nyckeltal.soliditet < 0:
            flaggor.append(RodFlagga(
                typ="negativ_soliditet",
                allvarlighet="kritisk",
                beskrivning=f"Negativ soliditet: {nyckeltal.soliditet:.1f}%",
                varde=nyckeltal.soliditet,
                rekommendation="Bolaget har mer skulder än tillgångar"
            ))
        elif nyckeltal.soliditet < 15:
            flaggor.append(RodFlagga(
                typ="lag_soliditet",
                allvarlighet="varning",
                beskrivning=f"Låg soliditet: {nyckeltal.soliditet:.1f}%",
                varde=nyckeltal.soliditet,
                rekommendation="Överväg att stärka det egna kapitalet"
            ))
    
    # 4. Negativ vinstmarginal
    if nyckeltal.vinstmarginal is not None and nyckeltal.vinstmarginal < -10:
        flaggor.append(RodFlagga(
            typ="negativ_vinstmarginal",
            allvarlighet="varning",
            beskrivning=f"Negativ vinstmarginal: {nyckeltal.vinstmarginal:.1f}%",
            varde=nyckeltal.vinstmarginal,
            rekommendation="Intäkterna täcker inte kostnaderna"
        ))
    
    return flaggor


def calculate_cagr(values: List[int], years: int) -> Optional[float]:
    """Beräkna CAGR (Compound Annual Growth Rate)."""
    if not values or len(values) < 2 or years < 1:
        return None
    start_value = values[-1]
    end_value = values[0]
    if start_value <= 0 or end_value <= 0:
        return None
    try:
        cagr = ((end_value / start_value) ** (1 / years) - 1) * 100
        return round(cagr, 2)
    except:
        return None


def linear_forecast(values: List[int]) -> Optional[int]:
    """Enkel linjär prognos för nästa värde."""
    clean = [v for v in values if v is not None]
    if len(clean) < 2:
        return None
    n = len(clean)
    x = list(range(n))
    y = clean[::-1]
    x_mean = sum(x) / n
    y_mean = sum(y) / n
    num = sum((x[i] - x_mean) * (y[i] - y_mean) for i in range(n))
    den = sum((x[i] - x_mean) ** 2 for i in range(n))
    if den == 0:
        return None
    slope = num / den
    intercept = y_mean - slope * x_mean
    return int(slope * n + intercept)


def fetch_full_arsredovisning(org_nummer: str, index: int = 0) -> Tuple[FullArsredovisning, bytes, bytes]:
    """Hämta och parsa KOMPLETT årsredovisning med alla detaljer."""
    clean_nr = clean_org_nummer(org_nummer)
    
    dok_data = make_api_request("POST", "/dokumentlista", {"identitetsbeteckning": clean_nr})
    dokument = dok_data.get("dokument", [])
    
    if not dokument:
        raise Exception("Inga årsredovisningar hittades")
    if index >= len(dokument):
        raise Exception(f"Index {index} finns inte.")
    
    dok = dokument[index]
    zip_bytes = download_document_bytes(dok.get("dokumentId"))
    
    xhtml_content = None
    with zipfile.ZipFile(BytesIO(zip_bytes)) as zf:
        for name in zf.namelist():
            if name.lower().endswith(('.xhtml', '.html', '.xml')):
                xhtml_content = zf.read(name).decode('utf-8')
                break
    
    if not xhtml_content:
        raise Exception("Ingen XHTML-fil hittades")
    
    parser = IXBRLParser(xhtml_content)
    metadata = parser.get_metadata()
    personer = parser.get_personer()
    nyckeltal = parser.get_nyckeltal('period0')
    nyckeltal_prev = parser.get_nyckeltal('period1')
    
    # Separera styrelse, revisorer, VD
    styrelse = []
    revisorer = []
    vd = None
    for p in personer:
        if 'revisor' in p.roll.lower():
            revisorer.append(p)
        elif 'vd' in p.roll.lower() or 'verkställande' in p.roll.lower():
            vd = p
        else:
            styrelse.append(p)
    
    # Flerårsdata
    flerarsdata = []
    for i in range(5):
        nt = parser.get_nyckeltal(f'period{i}')
        if nt.nettoomsattning or nt.arets_resultat:
            flerarsdata.append({'period': i, 'nyckeltal': asdict(nt)})
    
    # Förvaltningsberättelse (extraktion av textfält)
    forvaltningsberattelse = {}
    fb_patterns = ['VerksamhetenArt', 'VasentligaHandelser', 'ForvantadFramtidaUtveckling']
    for pattern in fb_patterns:
        for tag in parser.soup.find_all('ix:nonnumeric', {'name': lambda x: x and pattern in x}):
            text = tag.get_text(strip=True)
            if len(text) > 50:
                forvaltningsberattelse[pattern] = text[:2000]
                break
    
    full = FullArsredovisning(
        org_nummer=metadata['org_nummer'],
        foretag_namn=metadata['foretag_namn'],
        rakenskapsar_start=metadata['rakenskapsar_start'],
        rakenskapsar_slut=metadata['rakenskapsar_slut'],
        nyckeltal=nyckeltal,
        styrelse=styrelse,
        revisorer=revisorer,
        vd=vd,
        balansrakning=parser.get_balansrakning(),
        resultatrakning=parser.get_resultatrakning(),
        forvaltningsberattelse=forvaltningsberattelse,
        noter=parser.noter if hasattr(parser, 'noter') else {},
        flerarsdata=flerarsdata,
        roda_flaggor=analyze_roda_flaggor(parser, nyckeltal, nyckeltal_prev),
        metadata=metadata
    )
    
    return full, xhtml_content.encode('utf-8'), zip_bytes


# --- Input-modeller för nya verktyg ---

class TrendInput(BaseModel):
    model_config = ConfigDict(extra="forbid")
    org_nummer: str = Field(min_length=10, description="Organisationsnummer")
    max_years: int = Field(default=5, ge=2, le=10, description="Max antal år att analysera")


class NetworkInput(BaseModel):
    model_config = ConfigDict(extra="forbid")
    org_nummer: str = Field(min_length=10, description="Organisationsnummer")


class CompareInput(BaseModel):
    model_config = ConfigDict(extra="forbid")
    org_nummer_1: str = Field(min_length=10, description="Första organisationsnummer")
    org_nummer_2: str = Field(min_length=10, description="Andra organisationsnummer")


class FullAnalysisInput(BaseModel):
    model_config = ConfigDict(extra="forbid")
    org_nummer: str = Field(min_length=10, description="Organisationsnummer")
    index: int = Field(default=0, ge=0, description="Vilken årsredovisning (0=senaste)")
    format: str = Field(default="markdown", description="Format: markdown, json")


# --- VERKTYG 1: Fullständig årsredovisningsanalys ---

@mcp.tool()
def bolagsverket_analyze_full(params: FullAnalysisInput) -> str:
    """
    Komplett analys av årsredovisning - hämtar ZIP, extraherar, parsar ALL information.
    
    Autonomt verktyg som:
    1. Hämtar original ZIP från Bolagsverket
    2. Extraherar iXBRL/XHTML
    3. Parsar ALLA fält strukturerat
    
    Returnerar:
    - Nyckeltal med beräknade värden
    - Fullständig resultaträkning
    - Fullständig balansräkning  
    - Styrelse, VD, revisorer
    - Förvaltningsberättelse
    - Flerårsöversikt
    - Röda flaggor/varningar
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        full, _, _ = fetch_full_arsredovisning(result, params.index)
        
        if params.format == "json":
            return export_to_json(full)
        
        # Markdown-format
        def fmt(val):
            if val is None: return "-"
            if isinstance(val, int): return f"{val:,}"
            return str(val)
        
        lines = [
            f"# {full.foretag_namn}",
            f"## Årsredovisning {full.rakenskapsar_start} – {full.rakenskapsar_slut}",
            f"**Org.nr:** {format_org_nummer(full.org_nummer)}",
            "",
        ]
        
        # Röda flaggor
        if full.roda_flaggor:
            lines.extend(["## ⚠️ Varningar", ""])
            for f in full.roda_flaggor:
                icon = "🔴" if f.allvarlighet == "kritisk" else "🟡"
                lines.append(f"- {icon} **{f.beskrivning}**")
                if f.rekommendation:
                    lines.append(f"  - *{f.rekommendation}*")
            lines.append("")
        
        # Nyckeltal
        nt = full.nyckeltal
        lines.extend([
            "## Nyckeltal",
            "| Nyckeltal | Värde |",
            "|-----------|------:|",
            f"| Nettoomsättning | {fmt(nt.nettoomsattning)} SEK |",
            f"| Resultat efter fin. poster | {fmt(nt.resultat_efter_finansiella)} SEK |",
            f"| Årets resultat | {fmt(nt.arets_resultat)} SEK |",
            f"| Eget kapital | {fmt(nt.eget_kapital)} SEK |",
            f"| Soliditet | {fmt(nt.soliditet)} % |",
            f"| Vinstmarginal | {fmt(nt.vinstmarginal)} % |",
            f"| ROE | {fmt(nt.roe)} % |",
            f"| Antal anställda | {fmt(nt.antal_anstallda)} |",
            "",
        ])
        
        # Styrelse
        lines.extend(["## Styrelse och ledning", ""])
        if full.vd:
            lines.append(f"**VD:** {full.vd.fullnamn}")
        if full.styrelse:
            lines.extend(["", "| Namn | Roll |", "|------|------|"])
            for p in full.styrelse:
                lines.append(f"| {p.fullnamn} | {p.roll} |")
        if full.revisorer:
            lines.extend(["", "**Revisorer:**"])
            for p in full.revisorer:
                lines.append(f"- {p.fullnamn} ({p.roll})")
        lines.append("")
        
        # Balansräkning
        br = full.balansrakning
        lines.extend([
            "## Balansräkning",
            "### Tillgångar",
            f"- Immateriella: {fmt(br.get('tillgangar', {}).get('immateriella'))} SEK",
            f"- Materiella: {fmt(br.get('tillgangar', {}).get('materiella'))} SEK",
            f"- Finansiella: {fmt(br.get('tillgangar', {}).get('finansiella'))} SEK",
            f"- Varulager: {fmt(br.get('tillgangar', {}).get('varulager'))} SEK",
            f"- Kundfordringar: {fmt(br.get('tillgangar', {}).get('kundfordringar'))} SEK",
            f"- Kassa/Bank: {fmt(br.get('tillgangar', {}).get('kassa_bank'))} SEK",
            f"- **Summa tillgångar: {fmt(br.get('tillgangar', {}).get('summa_tillgangar'))} SEK**",
            "",
            "### Eget kapital och skulder",
            f"- Aktiekapital: {fmt(br.get('eget_kapital_skulder', {}).get('aktiekapital'))} SEK",
            f"- Balanserat resultat: {fmt(br.get('eget_kapital_skulder', {}).get('balanserat_resultat'))} SEK",
            f"- Summa eget kapital: {fmt(br.get('eget_kapital_skulder', {}).get('summa_eget_kapital'))} SEK",
            f"- Långfristiga skulder: {fmt(br.get('eget_kapital_skulder', {}).get('langfristiga_skulder'))} SEK",
            f"- Kortfristiga skulder: {fmt(br.get('eget_kapital_skulder', {}).get('kortfristiga_skulder'))} SEK",
            "",
        ])
        
        # Resultaträkning
        rr = full.resultatrakning
        lines.extend([
            "## Resultaträkning",
            f"- Nettoomsättning: {fmt(rr.get('nettoomsattning'))} SEK",
            f"- Övriga rörelseintäkter: {fmt(rr.get('ovriga_rorelseinktakter'))} SEK",
            f"- Rörelsekostnader: {fmt(rr.get('varor_handelsvaror'))} SEK",
            f"- Personalkostnader: {fmt(rr.get('personalkostnader'))} SEK",
            f"- Avskrivningar: {fmt(rr.get('avskrivningar'))} SEK",
            f"- **Rörelseresultat: {fmt(rr.get('rorelseresultat'))} SEK**",
            f"- Finansiella intäkter: {fmt(rr.get('finansiella_intakter'))} SEK",
            f"- Finansiella kostnader: {fmt(rr.get('finansiella_kostnader'))} SEK",
            f"- **Resultat efter fin. poster: {fmt(rr.get('resultat_efter_finansiella'))} SEK**",
            f"- Skatt: {fmt(rr.get('skatt'))} SEK",
            f"- **Årets resultat: {fmt(rr.get('arets_resultat'))} SEK**",
            "",
        ])
        
        # Flerårsöversikt
        if full.flerarsdata and len(full.flerarsdata) > 1:
            lines.extend(["## Flerårsöversikt", ""])
            headers = ["Nyckeltal"] + [f"År -{d['period']}" for d in full.flerarsdata]
            lines.append("| " + " | ".join(headers) + " |")
            lines.append("|" + "|".join(["---"] * len(headers)) + "|")
            
            for metric, label in [('nettoomsattning', 'Omsättning'), ('arets_resultat', 'Resultat'), ('eget_kapital', 'Eget kapital')]:
                row = [label]
                for d in full.flerarsdata:
                    row.append(fmt(d['nyckeltal'].get(metric)))
                lines.append("| " + " | ".join(row) + " |")
            lines.append("")
        
        # Förvaltningsberättelse
        if full.forvaltningsberattelse:
            lines.extend(["## Förvaltningsberättelse", ""])
            for key, text in full.forvaltningsberattelse.items():
                lines.extend([f"### {key}", text[:500] + "..." if len(text) > 500 else text, ""])
        
        return "\n".join(lines)
        
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


# --- VERKTYG 3: Koncernanalys / jämförelse ---

@mcp.tool()
def bolagsverket_compare(params: CompareInput) -> str:
    """
    Jämför två företag sida vid sida.
    
    Returnerar:
    - Jämförelse av nyckeltal
    - Gemensamma styrelseledamöter
    - Relativ prestanda
    """
    try:
        valid1, org1 = validate_org_nummer(params.org_nummer_1)
        valid2, org2 = validate_org_nummer(params.org_nummer_2)
        if not valid1:
            return handle_error(ErrorCode.INVALID_INPUT, org1)
        if not valid2:
            return handle_error(ErrorCode.INVALID_INPUT, org2)
        
        full1, _, _ = fetch_full_arsredovisning(org1, 0)
        full2, _, _ = fetch_full_arsredovisning(org2, 0)
        
        def fmt(val):
            if val is None: return "-"
            if isinstance(val, int): return f"{val:,}"
            if isinstance(val, float): return f"{val:.1f}"
            return str(val)
        
        lines = [
            "# Företagsjämförelse",
            "",
            f"| | {full1.foretag_namn} | {full2.foretag_namn} |",
            "|---|---:|---:|",
            f"| Org.nr | {format_org_nummer(full1.org_nummer)} | {format_org_nummer(full2.org_nummer)} |",
            f"| Omsättning | {fmt(full1.nyckeltal.nettoomsattning)} | {fmt(full2.nyckeltal.nettoomsattning)} |",
            f"| Resultat | {fmt(full1.nyckeltal.arets_resultat)} | {fmt(full2.nyckeltal.arets_resultat)} |",
            f"| Eget kapital | {fmt(full1.nyckeltal.eget_kapital)} | {fmt(full2.nyckeltal.eget_kapital)} |",
            f"| Soliditet | {fmt(full1.nyckeltal.soliditet)}% | {fmt(full2.nyckeltal.soliditet)}% |",
            f"| Vinstmarginal | {fmt(full1.nyckeltal.vinstmarginal)}% | {fmt(full2.nyckeltal.vinstmarginal)}% |",
            f"| Anställda | {fmt(full1.nyckeltal.antal_anstallda)} | {fmt(full2.nyckeltal.antal_anstallda)} |",
            "",
        ]
        
        # Hitta gemensamma personer
        personer1 = {p.fullnamn.lower() for p in full1.styrelse + full1.revisorer}
        if full1.vd: personer1.add(full1.vd.fullnamn.lower())
        personer2 = {p.fullnamn.lower() for p in full2.styrelse + full2.revisorer}
        if full2.vd: personer2.add(full2.vd.fullnamn.lower())
        
        gemensamma = personer1 & personer2
        if gemensamma:
            lines.extend(["## Gemensamma personer", ""])
            for namn in gemensamma:
                lines.append(f"- {namn.title()}")
            lines.append("")
        
        # Analys
        lines.extend(["## Analys", ""])
        
        if full1.nyckeltal.nettoomsattning and full2.nyckeltal.nettoomsattning:
            if full1.nyckeltal.nettoomsattning > full2.nyckeltal.nettoomsattning:
                lines.append(f"- {full1.foretag_namn} har {(full1.nyckeltal.nettoomsattning/full2.nyckeltal.nettoomsattning):.1f}x högre omsättning")
            else:
                lines.append(f"- {full2.foretag_namn} har {(full2.nyckeltal.nettoomsattning/full1.nyckeltal.nettoomsattning):.1f}x högre omsättning")
        
        if full1.nyckeltal.soliditet and full2.nyckeltal.soliditet:
            batter = full1.foretag_namn if full1.nyckeltal.soliditet > full2.nyckeltal.soliditet else full2.foretag_namn
            lines.append(f"- {batter} har bättre soliditet")
        
        return "\n".join(lines)
        
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


# --- VERKTYG 4: Trendvarningar (inkluderad i analyze_full) ---
# Röda flaggor är nu inkluderade i bolagsverket_analyze_full

@mcp.tool()
def bolagsverket_risk_check(params: OrgNummerInput) -> str:
    """
    Snabb riskanalys - identifierar röda flaggor och varningar.
    
    Kontrollerar:
    - Negativt eget kapital
    - Fallande omsättning (>20%)
    - Låg/negativ soliditet
    - Negativ vinstmarginal
    - Upprepade förluster
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        # Hämta både API-data och årsredovisningsdata
        company_info = fetch_company_info(result)
        full, _, _ = fetch_full_arsredovisning(result, 0)
        
        lines = [
            f"# Riskanalys: {full.foretag_namn}",
            f"**Org.nr:** {format_org_nummer(full.org_nummer)}",
            f"**Status:** {company_info.status}",
            "",
        ]
        
        # Samla alla flaggor - börja med de från årsredovisningen
        alla_flaggor = list(full.roda_flaggor)
        
        # Lägg till API-baserade flaggor
        
        # KRITISK: Pågående konkurs
        if company_info.pagaende_konkurs:
            alla_flaggor.append(RodFlagga(
                typ="pagaende_konkurs",
                allvarlighet="kritisk",
                beskrivning=f"⚠️ PÅGÅENDE KONKURS sedan {company_info.pagaende_konkurs.get('datum', 'okänt datum')[:10]}",
                varde=company_info.pagaende_konkurs,
                rekommendation="Företaget är under konkursförfarande. Undvik affärer."
            ))
        
        # KRITISK: Pågående likvidation
        if company_info.pagaende_likvidation:
            alla_flaggor.append(RodFlagga(
                typ="pagaende_likvidation",
                allvarlighet="kritisk",
                beskrivning=f"⚠️ PÅGÅENDE LIKVIDATION sedan {company_info.pagaende_likvidation.get('datum', 'okänt datum')[:10]}",
                varde=company_info.pagaende_likvidation,
                rekommendation="Företaget håller på att avvecklas."
            ))
        
        # KRITISK: Avregistrerat företag
        if company_info.status == "Avregistrerad":
            orsak = company_info.avregistreringsorsak or "Okänd orsak"
            alla_flaggor.append(RodFlagga(
                typ="avregistrerad",
                allvarlighet="kritisk",
                beskrivning=f"Företaget är AVREGISTRERAT ({orsak}) per {company_info.avregistreringsdatum or 'okänt datum'}",
                varde={"datum": company_info.avregistreringsdatum, "orsak": orsak},
                rekommendation="Företaget existerar inte längre som juridisk person."
            ))
        
        # INFO: Ej verksam organisation
        if company_info.verksam_organisation is False:
            alla_flaggor.append(RodFlagga(
                typ="ej_verksam",
                allvarlighet="info",
                beskrivning="Organisationen är markerad som EJ VERKSAM hos SCB",
                varde=None,
                rekommendation="Kan indikera vilande bolag eller felaktiga uppgifter."
            ))
        
        # INFO: Reklamspärr
        if company_info.reklamsparr is True:
            alla_flaggor.append(RodFlagga(
                typ="reklamsparr",
                allvarlighet="info",
                beskrivning="Företaget har reklamspärr",
                varde=None,
                rekommendation="Kan inte kontaktas via adresserad direktreklam."
            ))
        
        # INFO: Nyregistrerat bolag
        if company_info.registreringsdatum:
            try:
                reg_date = datetime.strptime(company_info.registreringsdatum[:10], "%Y-%m-%d")
                age_days = (datetime.now() - reg_date).days
                if age_days < 365 * 2:  # Mindre än 2 år
                    alla_flaggor.append(RodFlagga(
                        typ="nyregistrerat",
                        allvarlighet="info",
                        beskrivning=f"Nyregistrerat bolag ({age_days // 30} månader gammalt)",
                        varde=age_days,
                        rekommendation="Begränsad historik för bedömning."
                    ))
            except:
                pass
        
        # Sortera flaggor efter allvarlighet
        kritiska = [f for f in alla_flaggor if f.allvarlighet == "kritisk"]
        varningar = [f for f in alla_flaggor if f.allvarlighet == "varning"]
        info = [f for f in alla_flaggor if f.allvarlighet == "info"]
        
        if not alla_flaggor:
            lines.extend([
                "## ✅ Inga varningar",
                "",
                "Inga uppenbara röda flaggor identifierades.",
                "",
                "**Sammanfattning:**",
                f"- Soliditet: {full.nyckeltal.soliditet or '-'}%",
                f"- Vinstmarginal: {full.nyckeltal.vinstmarginal or '-'}%",
                f"- Eget kapital: {full.nyckeltal.eget_kapital:,} SEK" if full.nyckeltal.eget_kapital else "- Eget kapital: -",
            ])
        else:
            # Räkna risk-score
            risk_score = len(kritiska) * 3 + len(varningar) * 1
            risk_level = "🔴 HÖG" if risk_score >= 3 else "🟡 MEDIUM" if risk_score >= 1 else "🟢 LÅG"
            
            lines.extend([
                f"## Risknivå: {risk_level}",
                f"**{len(alla_flaggor)} observation(er):** {len(kritiska)} kritiska, {len(varningar)} varningar, {len(info)} info",
                "",
            ])
            
            if kritiska:
                lines.append("### 🔴 Kritiska")
                for f in kritiska:
                    lines.append(f"- **{f.beskrivning}**")
                    if f.rekommendation:
                        lines.append(f"  - _{f.rekommendation}_")
                lines.append("")
            
            if varningar:
                lines.append("### 🟡 Varningar")
                for f in varningar:
                    lines.append(f"- **{f.beskrivning}**")
                    if f.rekommendation:
                        lines.append(f"  - _{f.rekommendation}_")
                lines.append("")
            
            if info:
                lines.append("### ℹ️ Information")
                for f in info:
                    lines.append(f"- {f.beskrivning}")
                lines.append("")
        
        # Lägg till nyckeltal-sammanfattning
        lines.extend([
            "---",
            "## Nyckeltal (senaste år)",
            f"- Omsättning: {full.nyckeltal.nettoomsattning:,} SEK" if full.nyckeltal.nettoomsattning else "- Omsättning: -",
            f"- Resultat: {full.nyckeltal.arets_resultat:,} SEK" if full.nyckeltal.arets_resultat else "- Resultat: -",
            f"- Eget kapital: {full.nyckeltal.eget_kapital:,} SEK" if full.nyckeltal.eget_kapital else "- Eget kapital: -",
            f"- Soliditet: {full.nyckeltal.soliditet}%" if full.nyckeltal.soliditet else "- Soliditet: -",
        ])
        
        return "\n".join(lines)
        
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


# --- VERKTYG 7: Historisk tidsserie & prognoser ---

@mcp.tool()
def bolagsverket_trends(params: TrendInput) -> str:
    """
    Flerårsöversikt med trendanalys och prognoser.
    
    Hämtar upp till 10 års årsredovisningar och returnerar:
    - Nyckeltal per år
    - CAGR (årlig tillväxttakt)
    - Linjär prognos för nästa år
    - Trendvisualisering
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        clean_nr = clean_org_nummer(result)
        dok_data = make_api_request("POST", "/dokumentlista", {"identitetsbeteckning": clean_nr})
        dokument = dok_data.get("dokument", [])[:params.max_years]
        
        if len(dokument) < 2:
            return handle_error(ErrorCode.API_ERROR, "Minst 2 årsredovisningar krävs för trendanalys")
        
        # Hämta alla årsredovisningar
        data = []
        foretag_namn = None
        
        for i, dok in enumerate(dokument):
            try:
                zip_bytes = download_document_bytes(dok.get("dokumentId"))
                with zipfile.ZipFile(BytesIO(zip_bytes)) as zf:
                    for name in zf.namelist():
                        if name.lower().endswith(('.xhtml', '.html', '.xml')):
                            xhtml = zf.read(name).decode('utf-8')
                            parser = IXBRLParser(xhtml)
                            metadata = parser.get_metadata()
                            nyckeltal = parser.get_nyckeltal('period0')
                            
                            if not foretag_namn:
                                foretag_namn = metadata.get('foretag_namn', 'Okänt')
                            
                            data.append({
                                'period': metadata.get('rakenskapsar_slut', f'År {i}'),
                                'omsattning': nyckeltal.nettoomsattning,
                                'resultat': nyckeltal.arets_resultat,
                                'eget_kapital': nyckeltal.eget_kapital,
                                'anstallda': nyckeltal.antal_anstallda,
                            })
                            break
            except Exception as e:
                logger.warning(f"Kunde inte hämta årsredovisning {i}: {e}")
        
        if len(data) < 2:
            return handle_error(ErrorCode.API_ERROR, "Kunde inte hämta tillräckligt med data")
        
        def fmt(val):
            if val is None: return "-"
            if isinstance(val, int): return f"{val:,}"
            if isinstance(val, float): return f"{val:.1f}"
            return str(val)
        
        lines = [
            f"# Trendanalys: {foretag_namn}",
            f"**Analyserade år:** {len(data)}",
            "",
            "## Historik",
            "",
        ]
        
        # Tabell
        headers = ["Period", "Omsättning", "Resultat", "Eget kapital", "Anställda"]
        lines.append("| " + " | ".join(headers) + " |")
        lines.append("|" + "|".join(["---"] * len(headers)) + "|")
        
        for d in data:
            lines.append(f"| {d['period']} | {fmt(d['omsattning'])} | {fmt(d['resultat'])} | {fmt(d['eget_kapital'])} | {fmt(d['anstallda'])} |")
        
        lines.append("")
        
        # CAGR
        years = len(data) - 1
        oms_values = [d['omsattning'] for d in data if d['omsattning']]
        res_values = [d['resultat'] for d in data if d['resultat']]
        ek_values = [d['eget_kapital'] for d in data if d['eget_kapital']]
        
        lines.extend([
            "## Tillväxt (CAGR)",
            "",
        ])
        
        cagr_oms = calculate_cagr(oms_values, years)
        cagr_ek = calculate_cagr(ek_values, years)
        
        if cagr_oms:
            trend = "📈" if cagr_oms > 0 else "📉"
            lines.append(f"- Omsättning: {trend} **{cagr_oms:+.1f}%** per år")
        if cagr_ek:
            trend = "📈" if cagr_ek > 0 else "📉"
            lines.append(f"- Eget kapital: {trend} **{cagr_ek:+.1f}%** per år")
        
        lines.append("")
        
        # Prognos
        lines.extend([
            "## Prognos nästa år",
            "",
        ])
        
        prog_oms = linear_forecast(oms_values)
        prog_res = linear_forecast(res_values)
        prog_ek = linear_forecast(ek_values)
        
        if prog_oms:
            lines.append(f"- Omsättning: ~{prog_oms:,} SEK")
        if prog_res:
            lines.append(f"- Resultat: ~{prog_res:,} SEK")
        if prog_ek:
            lines.append(f"- Eget kapital: ~{prog_ek:,} SEK")
        
        lines.extend([
            "",
            "*Prognos baserad på linjär regression. Endast indikativ.*",
        ])
        
        return "\n".join(lines)
        
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


# --- VERKTYG 9: Styrelsenätverk ---

@mcp.tool()
def bolagsverket_network(params: NetworkInput) -> str:
    """
    Visa styrelsenätverk för ett företag.
    
    Returnerar:
    - Alla personer i styrelsen
    - VD och revisorer
    - Tips: Använd bolagsverket_compare för att hitta kopplingar mellan företag
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        full, _, _ = fetch_full_arsredovisning(result, 0)
        
        lines = [
            f"# Styrelsenätverk: {full.foretag_namn}",
            f"**Org.nr:** {format_org_nummer(full.org_nummer)}",
            "",
        ]
        
        if full.vd:
            lines.extend([
                "## VD",
                f"- **{full.vd.fullnamn}**",
                "",
            ])
        
        if full.styrelse:
            lines.extend([
                "## Styrelse",
                "",
            ])
            for p in full.styrelse:
                lines.append(f"- **{p.fullnamn}** - {p.roll}")
        
        if full.revisorer:
            lines.extend([
                "",
                "## Revisorer",
                "",
            ])
            for p in full.revisorer:
                lines.append(f"- **{p.fullnamn}** - {p.roll}")
        
        lines.extend([
            "",
            "---",
            "*Tips: Använd `bolagsverket_compare` för att hitta gemensamma styrelseledamöter mellan två företag.*",
        ])
        
        return "\n".join(lines)
        
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


# =============================================================================
# VERKTYG: Koncernanalys
# =============================================================================

class KoncernInput(BaseModel):
    model_config = ConfigDict(extra="forbid")
    org_nummer_lista: List[str] = Field(
        min_length=2,
        max_length=15,
        description="Lista med organisationsnummer för koncernbolag"
    )


@mcp.tool()
def bolagsverket_koncern(params: KoncernInput) -> str:
    """
    Analysera flera bolag som en koncern.
    
    Tar en lista med organisationsnummer och returnerar:
    - Sammanställning av nyckeltal för alla bolag
    - Summerade värden (omsättning, resultat, eget kapital)
    - Jämförelse mellan bolagen
    
    OBS: Summorna inkluderar eventuell koncernintern försäljning.
    """
    try:
        results = []
        errors = []
        alla_personer = {}  # namn -> lista av bolag
        
        for org_nr in params.org_nummer_lista:
            valid, clean = validate_org_nummer(org_nr)
            if not valid:
                errors.append(f"{org_nr}: Ogiltigt format")
                continue
            
            try:
                full, _, _ = fetch_full_arsredovisning(clean, 0)
                
                # Räkna varningar
                kritiska = len([f for f in full.roda_flaggor if f.allvarlighet == "kritisk"])
                varningar = len([f for f in full.roda_flaggor if f.allvarlighet == "varning"])
                
                results.append({
                    'org_nummer': format_org_nummer(clean),
                    'namn': full.foretag_namn,
                    'nyckeltal': full.nyckeltal,
                    'styrelse': full.styrelse,
                    'vd': full.vd,
                    'revisorer': full.revisorer,
                    'kritiska': kritiska,
                    'varningar': varningar,
                })
                
                # Samla personer för att hitta kopplingar
                for p in full.styrelse + full.revisorer + ([full.vd] if full.vd else []):
                    namn = p.fullnamn
                    if namn not in alla_personer:
                        alla_personer[namn] = []
                    alla_personer[namn].append({
                        'bolag': full.foretag_namn,
                        'roll': p.roll
                    })
                    
            except Exception as e:
                errors.append(f"{org_nr}: {str(e)[:50]}")
        
        if not results:
            return handle_error(ErrorCode.API_ERROR, f"Inga bolag kunde hämtas: {errors}")
        
        # Summera
        total_oms = sum(r['nyckeltal'].nettoomsattning or 0 for r in results)
        total_res = sum(r['nyckeltal'].arets_resultat or 0 for r in results)
        total_ek = sum(r['nyckeltal'].eget_kapital or 0 for r in results)
        total_anst = sum(r['nyckeltal'].antal_anstallda or 0 for r in results)
        total_kritiska = sum(r['kritiska'] for r in results)
        total_varningar = sum(r['varningar'] for r in results)
        
        def fmt(val):
            return f"{val:,}" if val else "-"
        
        lines = [
            "# Koncernöversikt",
            f"**Antal bolag:** {len(results)}",
            "",
        ]
        
        # Risk-sammanfattning
        if total_kritiska > 0 or total_varningar > 0:
            lines.extend([
                f"## ⚠️ Risksammanfattning",
                f"- 🔴 Kritiska varningar: {total_kritiska}",
                f"- 🟡 Varningar: {total_varningar}",
                "",
            ])
        
        lines.extend([
            "## Per bolag",
            "| Bolag | Omsättning | Resultat | EK | Anst. | Risk |",
            "|-------|----------:|---------:|---:|------:|:----:|",
        ])
        
        for r in results:
            nt = r['nyckeltal']
            namn = r['namn'][:22] + "..." if len(r['namn']) > 22 else r['namn']
            risk = "🔴" if r['kritiska'] > 0 else "🟡" if r['varningar'] > 0 else "✅"
            lines.append(f"| {namn} | {fmt(nt.nettoomsattning)} | {fmt(nt.arets_resultat)} | {fmt(nt.eget_kapital)} | {fmt(nt.antal_anstallda)} | {risk} |")
        
        lines.append(f"| **SUMMA** | **{fmt(total_oms)}** | **{fmt(total_res)}** | **{fmt(total_ek)}** | **{fmt(total_anst)}** | |")
        
        lines.extend([
            "",
            "## Koncerntotalt",
            f"- **Omsättning:** {fmt(total_oms)} SEK",
            f"- **Resultat:** {fmt(total_res)} SEK", 
            f"- **Eget kapital:** {fmt(total_ek)} SEK",
            f"- **Anställda:** {fmt(total_anst)}",
        ])
        
        # Beräkna aggregerad soliditet
        total_balans = sum(r['nyckeltal'].balansomslutning or 0 for r in results)
        if total_balans > 0:
            koncern_soliditet = round((total_ek / total_balans) * 100, 1)
            lines.append(f"- **Soliditet (aggregerad):** {koncern_soliditet}%")
        
        # Hitta gemensamma personer (sitter i flera bolag)
        gemensamma = {namn: bolag for namn, bolag in alla_personer.items() if len(bolag) > 1}
        
        if gemensamma:
            lines.extend([
                "",
                "## 🔗 Gemensamma nyckelpersoner",
                "",
            ])
            for namn, bolag_lista in gemensamma.items():
                lines.append(f"**{namn}**")
                for b in bolag_lista:
                    lines.append(f"  - {b['bolag']}: {b['roll']}")
                lines.append("")
        
        if errors:
            lines.extend(["", "## ⚠️ Kunde ej hämtas", ""])
            for err in errors:
                lines.append(f"- {err}")
        
        lines.extend([
            "",
            "---",
            "*OBS: Summorna kan inkludera koncernintern försäljning och elimineringar.*"
        ])
        
        return "\n".join(lines)
        
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


# =============================================================================
# VERKTYG: Export Word/PowerPoint
# =============================================================================

class ExportProInput(BaseModel):
    model_config = ConfigDict(extra="forbid")
    org_nummer: str = Field(min_length=10, description="Organisationsnummer")
    format: str = Field(default="word", description="'word' eller 'pptx'")


@mcp.tool()
def bolagsverket_export_pro(params: ExportProInput) -> str:
    """
    Exportera årsredovisning till Word eller PowerPoint.
    
    Kräver:
    - word: pip install python-docx
    - pptx: pip install python-pptx
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        # Hämta full årsredovisning för mer data
        full, _, _ = fetch_full_arsredovisning(result, 0)
        nt = full.nyckeltal
        
        clean_name = re.sub(r'[^\w\s-]', '', full.foretag_namn).strip().replace(' ', '_')
        year = full.rakenskapsar_slut[:4] if full.rakenskapsar_slut else "unknown"
        
        if params.format.lower() == "word":
            try:
                from docx import Document
                from docx.shared import Pt, Inches, RGBColor
                from docx.enum.text import WD_ALIGN_PARAGRAPH
            except ImportError:
                return "❌ python-docx ej installerat. Kör: pip install python-docx"
            
            doc = Document()
            
            # Titel
            title = doc.add_heading(full.foretag_namn, 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            subtitle = doc.add_paragraph(f"Årsredovisning {full.rakenskapsar_start} – {full.rakenskapsar_slut}")
            subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            doc.add_paragraph(f"Organisationsnummer: {format_org_nummer(result)}")
            doc.add_paragraph("")
            
            # Röda flaggor (om finns)
            if full.roda_flaggor:
                doc.add_heading("⚠️ Varningar", level=1)
                kritiska = [f for f in full.roda_flaggor if f.allvarlighet == "kritisk"]
                varningar = [f for f in full.roda_flaggor if f.allvarlighet == "varning"]
                
                if kritiska:
                    p = doc.add_paragraph()
                    run = p.add_run("KRITISKA:")
                    run.bold = True
                    run.font.color.rgb = RGBColor(192, 0, 0)
                    for f in kritiska:
                        doc.add_paragraph(f"• {f.beskrivning}", style='List Bullet')
                
                if varningar:
                    p = doc.add_paragraph()
                    run = p.add_run("Varningar:")
                    run.bold = True
                    for f in varningar:
                        doc.add_paragraph(f"• {f.beskrivning}", style='List Bullet')
                
                doc.add_paragraph("")
            
            # Nyckeltal
            doc.add_heading("Nyckeltal", level=1)
            table = doc.add_table(rows=1, cols=2)
            table.style = 'Table Grid'
            hdr = table.rows[0].cells
            hdr[0].text = 'Nyckeltal'
            hdr[1].text = 'Värde'
            
            for namn, val in [
                ('Nettoomsättning', f"{nt.nettoomsattning:,} SEK" if nt.nettoomsattning else '-'),
                ('Rörelseresultat', f"{full.resultatrakning.get('rorelseresultat', 0):,} SEK" if full.resultatrakning.get('rorelseresultat') else '-'),
                ('Resultat efter fin. poster', f"{nt.resultat_efter_finansiella:,} SEK" if nt.resultat_efter_finansiella else '-'),
                ('Årets resultat', f"{nt.arets_resultat:,} SEK" if nt.arets_resultat else '-'),
                ('Eget kapital', f"{nt.eget_kapital:,} SEK" if nt.eget_kapital else '-'),
                ('Balansomslutning', f"{nt.balansomslutning:,} SEK" if nt.balansomslutning else '-'),
                ('Soliditet', f"{nt.soliditet}%" if nt.soliditet else '-'),
                ('Vinstmarginal', f"{nt.vinstmarginal}%" if nt.vinstmarginal else '-'),
                ('ROE', f"{nt.roe}%" if nt.roe else '-'),
                ('Antal anställda', str(nt.antal_anstallda) if nt.antal_anstallda else '-'),
            ]:
                row = table.add_row().cells
                row[0].text = namn
                row[1].text = val
            
            doc.add_paragraph("")
            
            # Styrelse och ledning
            doc.add_heading("Styrelse och ledning", level=1)
            
            if full.vd:
                doc.add_paragraph(f"VD: {full.vd.fullnamn}")
            
            if full.styrelse:
                doc.add_paragraph("Styrelse:")
                for p in full.styrelse:
                    doc.add_paragraph(f"• {p.fullnamn} - {p.roll}", style='List Bullet')
            
            if full.revisorer:
                doc.add_paragraph("Revisorer:")
                for p in full.revisorer:
                    doc.add_paragraph(f"• {p.fullnamn} ({p.roll})", style='List Bullet')
            
            # Förvaltningsberättelse (utdrag)
            if full.forvaltningsberattelse:
                doc.add_heading("Förvaltningsberättelse (utdrag)", level=1)
                for key, text in full.forvaltningsberattelse.items():
                    if text and len(text) > 20:
                        # Begränsa längden
                        text_short = text[:500] + "..." if len(text) > 500 else text
                        doc.add_paragraph(text_short)
                        break
            
            filename = f"{clean_name}_{year}_rapport.docx"
            filepath = ensure_output_dir() / filename
            doc.save(filepath)
            
            return f"✅ Word-dokument sparat:\n📁 {filepath}"
            
        elif params.format.lower() == "pptx":
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.enum.text import PP_ALIGN
                from pptx.dml.color import RGBColor
            except ImportError:
                return "❌ python-pptx ej installerat. Kör: pip install python-pptx"
            
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # Slide 1: Titel
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
            tf = title.text_frame
            p = tf.paragraphs[0]
            p.text = full.foretag_namn
            p.font.size = Pt(48)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
            sub = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(0.8))
            tf = sub.text_frame
            p = tf.paragraphs[0]
            p.text = f"Årsredovisning {full.rakenskapsar_start} – {full.rakenskapsar_slut}"
            p.font.size = Pt(28)
            p.alignment = PP_ALIGN.CENTER
            
            org_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12), Inches(0.5))
            tf = org_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Org.nr: {format_org_nummer(result)}"
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.CENTER
            
            # Slide 2: Nyckeltal
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
            tf = title.text_frame
            p = tf.paragraphs[0]
            p.text = "Nyckeltal"
            p.font.size = Pt(40)
            p.font.bold = True
            
            # Vänster kolumn
            left = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6), Inches(5.5))
            tf = left.text_frame
            
            items_left = [
                f"Omsättning: {nt.nettoomsattning:,} SEK" if nt.nettoomsattning else "Omsättning: -",
                f"Resultat: {nt.arets_resultat:,} SEK" if nt.arets_resultat else "Resultat: -",
                f"Eget kapital: {nt.eget_kapital:,} SEK" if nt.eget_kapital else "Eget kapital: -",
            ]
            
            for i, item in enumerate(items_left):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(28)
                p.space_after = Pt(20)
            
            # Höger kolumn
            right = slide.shapes.add_textbox(Inches(6.5), Inches(1.3), Inches(6), Inches(5.5))
            tf = right.text_frame
            
            items_right = [
                f"Soliditet: {nt.soliditet}%" if nt.soliditet else "Soliditet: -",
                f"Vinstmarginal: {nt.vinstmarginal}%" if nt.vinstmarginal else "Vinstmarginal: -",
                f"Anställda: {nt.antal_anstallda}" if nt.antal_anstallda else "Anställda: -",
            ]
            
            for i, item in enumerate(items_right):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(28)
                p.space_after = Pt(20)
            
            # Slide 3: Varningar (om finns)
            if full.roda_flaggor:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
                tf = title.text_frame
                p = tf.paragraphs[0]
                p.text = "⚠️ Varningar"
                p.font.size = Pt(40)
                p.font.bold = True
                
                content = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
                tf = content.text_frame
                
                for i, f in enumerate(full.roda_flaggor[:6]):  # Max 6 på en slide
                    icon = "🔴" if f.allvarlighet == "kritisk" else "🟡"
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"{icon} {f.beskrivning}"
                    p.font.size = Pt(22)
                    p.space_after = Pt(15)
            
            # Slide 4: Styrelse
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
            tf = title.text_frame
            p = tf.paragraphs[0]
            p.text = "Styrelse och ledning"
            p.font.size = Pt(40)
            p.font.bold = True
            
            content = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
            tf = content.text_frame
            
            personer_text = []
            if full.vd:
                personer_text.append(f"VD: {full.vd.fullnamn}")
            for p in full.styrelse[:5]:
                personer_text.append(f"{p.roll}: {p.fullnamn}")
            
            for i, text in enumerate(personer_text):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = text
                p.font.size = Pt(24)
                p.space_after = Pt(15)
            
            filename = f"{clean_name}_{year}_presentation.pptx"
            filepath = ensure_output_dir() / filename
            prs.save(filepath)
            
            return f"✅ PowerPoint sparat:\n📁 {filepath}"
        
        else:
            return f"❌ Okänt format: {params.format}. Använd 'word' eller 'pptx'."
            
    except Exception as e:
        return handle_error(ErrorCode.EXPORT_ERROR, str(e))


# =============================================================================
# NYA MCP-TOOLS FÖR v5.1.0
# =============================================================================

class TaxonomiInput(OrgNummerMixin):
    """Input för taxonomi-verktyg."""
    pass

@mcp.tool()
def bolagsverket_taxonomy_info(params: TaxonomiInput) -> str:
    """
    NYTT I v5.1: Visa taxonomiversion och varningar för ett företags årsredovisning.
    
    Detekterar automatiskt vilken taxonomiversion som används och varnar för
    arkiverade versioner som ej längre stöds av Bolagsverket.
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        # Hämta årsredovisning
        arsredovisning, xhtml, _ = fetch_and_parse_arsredovisning(result, 0)
        parser = IXBRLParser(xhtml.decode('utf-8'))
        
        # Hämta taxonomi-info
        taxonomier = parser.get_taxonomi_info()
        varningar = parser.get_taxonomi_varningar()
        ar_koncern = parser.ar_koncernredovisning()
        
        output = [
            f"📋 TAXONOMIANALYS: {arsredovisning.foretag_namn}",
            f"Org.nr: {format_org_nummer(result)}",
            f"Räkenskapsår: {arsredovisning.rakenskapsar_slut}",
            ""
        ]
        
        if taxonomier:
            output.append("📚 DETEKTERADE TAXONOMIER:")
            for tax in taxonomier:
                status = "⚠️ ARKIVERAD" if tax.ar_arkiverad else "✅ Aktuell"
                output.append(f"  • {tax.typ} version {tax.version} [{status}]")
        else:
            output.append("❓ Ingen taxonomi kunde detekteras")
        
        output.append("")
        output.append(f"📊 Dokumenttyp: {'Koncernredovisning (K3K)' if ar_koncern else 'Årsredovisning'}")
        
        if varningar:
            output.append("")
            output.append("⚠️ VARNINGAR:")
            for v in varningar:
                output.append(f"  {v}")
        
        return "\n".join(output)
        
    except Exception as e:
        return handle_error(ErrorCode.PARSE_ERROR, str(e))


@mcp.tool()
def bolagsverket_bas_mapping(params: TaxonomiInput) -> str:
    """
    NYTT I v5.1: Visa BAS-kontomappning för ett företags årsredovisning.
    
    Mappar taxonomi-begrepp till BAS-kontoplanen för att underlätta
    integrering med bokföringssystem.
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        arsredovisning, xhtml, _ = fetch_and_parse_arsredovisning(result, 0)
        parser = IXBRLParser(xhtml.decode('utf-8'))
        
        mappings = parser.get_bas_mappning()
        
        output = [
            f"📒 BAS-KONTOMAPPNING: {arsredovisning.foretag_namn}",
            f"Org.nr: {format_org_nummer(result)}",
            ""
        ]
        
        if mappings:
            output.append(f"Hittade {len(mappings)} begrepp med BAS-mappning:")
            output.append("")
            
            # Gruppera efter kontogrupp
            for begrepp, data in sorted(mappings.items(), key=lambda x: x[1]['bas_konto']):
                varde = data['varde']
                bas = data['bas_konto']
                output.append(f"  {bas}: {begrepp} = {varde:,} kr")
        else:
            output.append("Inga begrepp kunde mappas till BAS-kontoplanen.")
        
        return "\n".join(output)
        
    except Exception as e:
        return handle_error(ErrorCode.PARSE_ERROR, str(e))


@mcp.tool()
def bolagsverket_koncern_extended(params: TaxonomiInput) -> str:
    """
    NYTT I v5.1: Utökad koncernanalys med K3K-nyckeltal.
    
    Extraherar koncernspecifika nyckeltal som:
    - Koncernens nettoomsättning
    - Koncernens resultat
    - Minoritetsandelar
    - Goodwill
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        arsredovisning, xhtml, _ = fetch_and_parse_arsredovisning(result, 0)
        parser = IXBRLParser(xhtml.decode('utf-8'))
        
        output = [
            f"🏢 KONCERNANALYS: {arsredovisning.foretag_namn}",
            f"Org.nr: {format_org_nummer(result)}",
            ""
        ]
        
        if not parser.ar_koncernredovisning():
            output.append("ℹ️ Detta är INTE en koncernredovisning.")
            output.append("Dokumentet innehåller inga koncernbegrepp.")
            return "\n".join(output)
        
        output.append("✅ Detta är en koncernredovisning (K3K)")
        output.append("")
        
        koncern = parser.get_koncern_nyckeltal('period0')
        if koncern:
            output.append("📊 KONCERNNYCKELTAL:")
            if koncern.koncern_nettoomsattning:
                output.append(f"  Nettoomsättning: {koncern.koncern_nettoomsattning:,} kr")
            if koncern.koncern_rorelseresultat:
                output.append(f"  Rörelseresultat: {koncern.koncern_rorelseresultat:,} kr")
            if koncern.koncern_arets_resultat:
                output.append(f"  Årets resultat: {koncern.koncern_arets_resultat:,} kr")
            if koncern.koncern_eget_kapital:
                output.append(f"  Eget kapital: {koncern.koncern_eget_kapital:,} kr")
            if koncern.minoritetsandel:
                output.append(f"  Minoritetsandel: {koncern.minoritetsandel:,} kr")
            if koncern.goodwill:
                output.append(f"  Goodwill: {koncern.goodwill:,} kr")
            if koncern.koncern_soliditet:
                output.append(f"  Soliditet: {koncern.koncern_soliditet:.1f}%")
        else:
            output.append("Inga koncernnyckeltal kunde extraheras.")
        
        return "\n".join(output)
        
    except Exception as e:
        return handle_error(ErrorCode.PARSE_ERROR, str(e))


@mcp.tool()
def bolagsverket_revision(params: TaxonomiInput) -> str:
    """
    NYTT I v5.1: Extrahera och visa revisionsberättelse.
    
    Visar strukturerad information om:
    - Revisor och revisionsbolag
    - Uttalanden
    - Eventuella anmärkningar
    - Om det är en "ren" revisionsberättelse
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        arsredovisning, xhtml, _ = fetch_and_parse_arsredovisning(result, 0)
        parser = IXBRLParser(xhtml.decode('utf-8'))
        
        rb = parser.get_revisionsberattelse()
        
        output = [
            f"📝 REVISIONSBERÄTTELSE: {arsredovisning.foretag_namn}",
            f"Org.nr: {format_org_nummer(result)}",
            ""
        ]
        
        if not rb:
            output.append("ℹ️ Ingen revisionsberättelse hittades i dokumentet.")
            output.append("(Bolaget kanske saknar revisionsplikt)")
            return "\n".join(output)
        
        output.append(f"📋 Typ: {rb.typ.upper()}")
        output.append(f"{'✅ REN' if rb.ar_ren else '⚠️ INNEHÅLLER ANMÄRKNINGAR'}")
        output.append("")
        
        if rb.revisor_namn:
            output.append(f"👤 Revisor: {rb.revisor_namn}")
        if rb.revisor_titel:
            output.append(f"   Titel: {rb.revisor_titel}")
        if rb.revisionsbolag:
            output.append(f"   Bolag: {rb.revisionsbolag}")
        
        if rb.datum or rb.ort:
            output.append(f"📅 {rb.ort or ''} {rb.datum or ''}")
        
        if rb.anmarkningar:
            output.append("")
            output.append("⚠️ ANMÄRKNINGAR:")
            for anm in rb.anmarkningar:
                output.append(f"  • {anm[:200]}...")
        
        return "\n".join(output)
        
    except Exception as e:
        return handle_error(ErrorCode.PARSE_ERROR, str(e))


@mcp.tool()
def bolagsverket_faststallelse(params: TaxonomiInput) -> str:
    """
    NYTT I v5.1: Extrahera fastställelseintyg.
    
    Visar information om:
    - Årsstämma och fastställelsedatum
    - Resultatdisposition (utdelning, balansering)
    - Undertecknare
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        arsredovisning, xhtml, _ = fetch_and_parse_arsredovisning(result, 0)
        parser = IXBRLParser(xhtml.decode('utf-8'))
        
        fi = parser.get_faststallelseintyg()
        
        output = [
            f"📜 FASTSTÄLLELSEINTYG: {arsredovisning.foretag_namn}",
            f"Org.nr: {format_org_nummer(result)}",
            ""
        ]
        
        if not fi:
            output.append("ℹ️ Inget fastställelseintyg hittades i dokumentet.")
            return "\n".join(output)
        
        if fi.arsstamma_datum:
            output.append(f"📅 Årsstämma: {fi.arsstamma_datum}")
        if fi.intygsdatum:
            output.append(f"📝 Intyg upprättat: {fi.intygsdatum}")
        
        output.append("")
        output.append("📊 RESULTATDISPOSITION:")
        if fi.utdelning_totalt:
            output.append(f"  Utdelning: {fi.utdelning_totalt:,} kr")
        if fi.utdelning_per_aktie:
            output.append(f"  Per aktie: {fi.utdelning_per_aktie} kr")
        if fi.balanseras_i_ny_rakning:
            output.append(f"  Balanseras: {fi.balanseras_i_ny_rakning:,} kr")
        
        if fi.undertecknare:
            output.append("")
            output.append("✍️ UNDERTECKNARE:")
            for namn in fi.undertecknare:
                output.append(f"  • {namn}")
        
        return "\n".join(output)
        
    except Exception as e:
        return handle_error(ErrorCode.PARSE_ERROR, str(e))


class CacheInput(BaseModel):
    """Input för cache-verktyg."""
    action: str = Field(default="stats", description="Åtgärd: stats, clear, clear_expired")

@mcp.tool()
def bolagsverket_cache(params: CacheInput) -> str:
    """
    NYTT I v5.1: Hantera cache för API-anrop.
    
    Åtgärder:
    - stats: Visa cache-statistik
    - clear: Rensa hela cachen
    - clear_expired: Rensa utgångna entries
    """
    try:
        if params.action == "stats":
            stats = cache_manager.get_stats()
            output = [
                "📦 CACHE-STATISTIK",
                "",
                f"Totalt: {stats['total_entries']} entries",
                f"Utgångna: {stats['expired_entries']} entries",
                f"Storlek: {stats['db_size_bytes'] / 1024:.1f} KB",
                ""
            ]
            
            if stats.get('categories'):
                output.append("Per kategori:")
                for cat, data in stats['categories'].items():
                    output.append(f"  {cat}: {data['count']} entries, {data['hits']} träffar")
            
            return "\n".join(output)
        
        elif params.action == "clear":
            count = cache_manager.clear_all()
            return f"✅ Rensade {count} cache-entries"
        
        elif params.action == "clear_expired":
            count = cache_manager.clear_expired()
            return f"✅ Rensade {count} utgångna cache-entries"
        
        else:
            return f"❌ Okänd åtgärd: {params.action}. Använd: stats, clear, clear_expired"
            
    except Exception as e:
        return handle_error(ErrorCode.UNKNOWN_ERROR, str(e))


@mcp.tool()
def bolagsverket_utokad_info(params: TaxonomiInput) -> str:
    """
    NYTT I v5.1: Visa utökad information från extension taxonomy.
    
    Extraherar:
    - Odefinierade begrepp (begrepp utanför standardtaxonomin)
    - Ändrade rubriker
    - Notkopplingar
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        arsredovisning, xhtml, _ = fetch_and_parse_arsredovisning(result, 0)
        parser = IXBRLParser(xhtml.decode('utf-8'))
        
        ui = parser.get_utokad_information()
        
        output = [
            f"🔍 UTÖKAD INFORMATION: {arsredovisning.foretag_namn}",
            f"Org.nr: {format_org_nummer(result)}",
            "",
            f"{'✅ Fullständigt taggad' if ui.ar_fullstandigt_taggad else '⚠️ Innehåller otaggad information'}",
            ""
        ]
        
        if ui.odefinierade_begrepp:
            output.append(f"📝 ODEFINIERADE BEGREPP ({len(ui.odefinierade_begrepp)} st):")
            for b in ui.odefinierade_begrepp[:10]:
                output.append(f"  • {b['namn']}: {b['varde']:,} kr" if b['varde'] else f"  • {b['namn']}")
        
        if ui.andrade_rubriker:
            output.append("")
            output.append(f"✏️ ÄNDRADE RUBRIKER ({len(ui.andrade_rubriker)} st):")
            for r in ui.andrade_rubriker[:10]:
                output.append(f"  • '{r['ursprunglig']}' → '{r['ny']}'")
        
        if ui.notkopplingar:
            output.append("")
            output.append(f"🔗 NOTKOPPLINGAR ({len(ui.notkopplingar)} st):")
            for n in ui.notkopplingar[:10]:
                output.append(f"  • Not {n['not_nummer']}")
        
        if not (ui.odefinierade_begrepp or ui.andrade_rubriker or ui.notkopplingar):
            output.append("Ingen utökad information hittades.")
        
        return "\n".join(output)
        
    except Exception as e:
        return handle_error(ErrorCode.PARSE_ERROR, str(e))


# =============================================================================
# Huvudprogram
# =============================================================================

def main():
    """
    Kör MCP-servern med vald transport.

    Användning:
        python bolagsverket_mcp_server.py           # stdio (Claude Desktop)
        python bolagsverket_mcp_server.py --http    # HTTP/SSE (remote clients)
        python bolagsverket_mcp_server.py --http --port 8080
    
    Miljövariabler:
        PORT                      - Port för HTTP-server (Render sätter denna)
        RENDER                    - Sätts automatiskt av Render
        BOLAGSVERKET_CLIENT_ID    - API Client ID
        BOLAGSVERKET_CLIENT_SECRET - API Client Secret
    """
    import argparse

    parser = argparse.ArgumentParser(description="Bolagsverket MCP Server v5.1.0")
    parser.add_argument(
        "--http",
        action="store_true",
        help="Kör med HTTP/SSE transport (för remote clients)"
    )
    parser.add_argument(
        "--port",
        type=int,
        default=None,
        help="Port för HTTP-server (default: $PORT eller 8000)"
    )
    parser.add_argument(
        "--host",
        type=str,
        default="0.0.0.0",
        help="Host för HTTP-server (default: 0.0.0.0)"
    )

    args = parser.parse_args()
    
    # Render sätter PORT miljövariabel
    port = args.port or int(os.environ.get("PORT", 8000))

    if args.http or os.environ.get("RENDER"):
        logger.info(f"Startar Bolagsverket MCP Server v5.1.0 (HTTP/SSE) på {args.host}:{port}...")
        mcp.run(transport="sse", host=args.host, port=port)
    else:
        logger.info("Startar Bolagsverket MCP Server v5.1.0 (stdio)...")
        mcp.run(transport="stdio")


if __name__ == "__main__":
    main()
